from __future__ import annotations

import csv
import hashlib
import os
import re
from pathlib import Path


ROOT = Path.cwd()
OUT_DIR = ROOT / "reports" / "选必一_哲学宝典式重建_2026-05-16" / "11_strict_final_rebuild_2026-05-23"
INVENTORY = OUT_DIR / "00_source_inventory.csv"
TEXT_DIR = OUT_DIR / "01_extracted_text"
LOG_CSV = OUT_DIR / "01_extraction_log.csv"
LOG_MD = OUT_DIR / "01_extraction_log.md"


def stable_id(path: str) -> str:
    return hashlib.sha1(path.encode("utf-8")).hexdigest()[:12]


def safe_name(name: str) -> str:
    name = re.sub(r"[\\/:*?\"<>|]+", "_", name)
    name = re.sub(r"\s+", "_", name).strip("_")
    return name[:120] or "source"


def extract_docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("[TBL] " + " | ".join(cells))
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    import fitz

    doc = fitz.open(str(path))
    parts: list[str] = []
    for idx, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        parts.append(f"--- page {idx} ---\n{text}")
    doc.close()
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    pres = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(pres.slides, start=1):
        parts.append(f"--- slide {idx} ---")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append("[TBL] " + " | ".join(cells))
    return "\n".join(parts)


def extract_doc_binary(path: Path) -> str:
    raw = path.read_bytes()
    texts: list[str] = []
    for encoding in ("utf-16-le", "gb18030", "latin-1"):
        try:
            decoded = raw.decode(encoding, errors="ignore")
        except Exception:
            continue
        keep: list[str] = []
        buf: list[str] = []
        for ch in decoded:
            cp = ord(ch)
            if (
                0x4E00 <= cp <= 0x9FFF
                or 0x3400 <= cp <= 0x4DBF
                or 0x3000 <= cp <= 0x303F
                or 0xFF00 <= cp <= 0xFFEF
                or ch in "\n\t"
                or 0x20 <= cp <= 0x7E
            ):
                buf.append(ch)
            else:
                if buf:
                    run = "".join(buf).strip()
                    if len(run) >= 2:
                        keep.append(run)
                    buf = []
        if buf:
            run = "".join(buf).strip()
            if len(run) >= 2:
                keep.append(run)
        text = "\n".join(keep)
        if len(text) > max((len(t) for t in texts), default=0):
            texts.append(text)
    return max(texts, key=len) if texts else ""


def extract_plain(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def quality_flag(text: str, ext: str) -> str:
    visible = re.sub(r"\s+", "", text)
    cjk = re.findall(r"[\u4e00-\u9fff]", text)
    if ext in {".png", ".jpg", ".jpeg"}:
        return "image_only"
    if len(visible) == 0:
        return "empty"
    if len(cjk) < 80 and ext in {".pdf", ".pptx", ".ppt", ".doc"}:
        return "needs_ocr_or_render_review"
    if len(visible) < 120:
        return "thin_text_review"
    return "usable_text"


def extract_one(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    if ext == ".docx":
        return extract_docx(path), "python-docx"
    if ext == ".pdf":
        return extract_pdf(path), "pymupdf_text"
    if ext == ".pptx":
        return extract_pptx(path), "python-pptx"
    if ext == ".doc":
        return extract_doc_binary(path), "binary_text_scan"
    if ext in {".txt", ".md", ".csv", ".json"}:
        return extract_plain(path), "plain_text"
    if ext in {".png", ".jpg", ".jpeg"}:
        return "", "image_unextracted"
    return "", "unsupported"


def main() -> None:
    TEXT_DIR.mkdir(parents=True, exist_ok=True)
    rows: list[dict[str, str]] = []
    with INVENTORY.open("r", encoding="utf-8-sig", newline="") as f:
        for item in csv.DictReader(f):
            source_path = Path(item["absolute_path"])
            sid = stable_id(item["absolute_path"])
            out_name = f"{sid}_{safe_name(source_path.stem)}.txt"
            out_path = TEXT_DIR / out_name
            status = "ok"
            error = ""
            text = ""
            method = ""
            try:
                text, method = extract_one(source_path)
                if method in {"unsupported", "image_unextracted"}:
                    status = method
                out_path.write_text(text, encoding="utf-8")
            except Exception as exc:
                status = "error"
                method = method or "unknown"
                error = repr(exc)
                out_path.write_text("", encoding="utf-8")
            flag = quality_flag(text, source_path.suffix.lower())
            row = {
                **item,
                "source_id": sid,
                "text_path": str(out_path),
                "method": method,
                "status": status,
                "quality_flag": flag,
                "char_count": str(len(text)),
                "cjk_count": str(len(re.findall(r"[\u4e00-\u9fff]", text))),
                "error": error,
            }
            rows.append(row)

    fields = list(rows[0].keys()) if rows else []
    with LOG_CSV.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fields)
        writer.writeheader()
        writer.writerows(rows)

    counts: dict[str, int] = {}
    quality: dict[str, int] = {}
    for row in rows:
        counts[row["status"]] = counts.get(row["status"], 0) + 1
        quality[row["quality_flag"]] = quality.get(row["quality_flag"], 0) + 1
    lines = [
        "# 选必一严格最终版源文件文本抽取日志",
        "",
        f"- 源文件：{len(rows)}",
        f"- 成功/状态：{counts}",
        f"- 文本质量：{quality}",
        "",
        "## 需复核文件",
        "",
        "| 质量标记 | 角色 | 文件 |",
        "|---|---|---|",
    ]
    for row in rows:
        if row["quality_flag"] != "usable_text":
            lines.append(f"| {row['quality_flag']} | {row['role']} | `{row['absolute_path']}` |")
    LOG_MD.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")

    print(f"rows={len(rows)}")
    print(LOG_CSV)
    print(LOG_MD)
    print(f"quality={quality}")


if __name__ == "__main__":
    main()
