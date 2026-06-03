from __future__ import annotations

import csv
import re
import subprocess
import sys
from pathlib import Path

sys.path.insert(0, "/tmp/xuanbiyi_review237_deps")

import fitz  # type: ignore
from docx import Document
from pptx import Presentation  # type: ignore


RUN_DIR = Path("/Users/wanglifei/GaokaoPolitics/beijing-politics-sync/reports/选必一_哲学宝典式重建_2026-05-16/20_review237_verify_final_20260528")
OUT_DIR = RUN_DIR / "source_texts"
F_OCR_DIR = Path("/Users/wanglifei/GaokaoPolitics/beijing-politics-sync/reports/选必一_哲学宝典式重建_2026-05-16/19_error_report_patch_20260527/f_class_ocr")

SUITE_DIRS = {
    "2024海淀二模": [Path("/Users/wanglifei/Desktop/2024模拟题/海淀二模")],
    "2026海淀期中": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026海淀期中")],
    "2025丰台期末": [Path("/Users/wanglifei/Desktop/2025模拟题/2025各区期末/2025丰台期末")],
    "2024西城二模": [Path("/Users/wanglifei/Desktop/2024模拟题/西城二模")],
    "2026东城一模": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区一模/2026东城一模")],
    "2026海淀一模": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区一模/2026海淀一模")],
    "2026丰台期末": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026丰台期末")],
    "2025东城期末": [Path("/Users/wanglifei/Desktop/2025模拟题/2025各区期末/2025东城期末")],
    "2026朝阳期末": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026朝阳期末")],
    "2026西城二模": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区二模/2026西城二模")],
    "2026顺义二模": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区二模/2026顺义二模")],
    "2025西城二模": [Path("/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025西城二模")],
    "2024朝阳一模": [Path("/Users/wanglifei/Desktop/2024模拟题/2024朝阳一模")],
    "2026西城期末": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026西城期末")],
    "2025顺义一模": [Path("/Users/wanglifei/Desktop/2025模拟题/2025各区一模/2025顺义一模")],
    "2026海淀二模": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区二模/2026海淀二模")],
    "2026顺义一模": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区一模/2026顺义一模")],
    "2026门头沟一模": [Path("/Users/wanglifei/Desktop/2026模拟题/2026各区一模/2026门头沟一模")],
    "2025西城一模": [Path("/Users/wanglifei/Desktop/2025模拟题/2025各区一模/2025西城一模")],
}

F_OCR_MAP = {
    "2026海淀一模": ["2026海淀一模_试卷.md", "2026海淀一模_细则.md"],
    "2026丰台期末": ["2026丰台期末_试卷.md", "2026丰台期末_细则.md"],
    "2026朝阳期末": ["2026朝阳期末_试卷.md", "2026朝阳期末_细则.md"],
    "2026顺义二模": ["2026顺义二模_试卷.md"],
    "2026西城期末": ["2026西城期末_参考答案.md", "2026西城期末_细则.md"],
}


def clean(text: str) -> str:
    text = text.replace("\x00", " ")
    return re.sub(r"[ \t\r\f\v]+", " ", text)


def docx_text(path: Path) -> str:
    doc = Document(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return clean("\n".join(parts))


def pptx_text(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        parts.append(f"--- slide {idx} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return clean("\n".join(parts))


def pdf_text(path: Path) -> str:
    doc = fitz.open(str(path))
    parts: list[str] = []
    for idx, page in enumerate(doc, 1):
        txt = page.get_text("text")
        if txt.strip():
            parts.append(f"--- page {idx} ---\n{txt}")
    return clean("\n".join(parts))


def doc_text(path: Path) -> str:
    try:
        out = subprocess.check_output(["textutil", "-convert", "txt", "-stdout", str(path)], stderr=subprocess.DEVNULL)
        return clean(out.decode("utf-8", errors="ignore"))
    except Exception:
        return ""


def extract(path: Path) -> str:
    suffix = path.suffix.lower()
    if path.name.startswith("~$"):
        return ""
    if suffix == ".docx":
        return docx_text(path)
    if suffix == ".pptx":
        return pptx_text(path)
    if suffix == ".pdf":
        return pdf_text(path)
    if suffix == ".doc":
        return doc_text(path)
    if suffix in {".txt", ".md", ".rtf"}:
        return clean(path.read_text(encoding="utf-8", errors="ignore"))
    return ""


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    rows: list[dict[str, str]] = []
    for suite, roots in SUITE_DIRS.items():
        suite_parts: list[str] = []
        for root in roots:
            if not root.exists():
                rows.append({"suite": suite, "file": str(root), "chars": "0", "status": "missing_root"})
                continue
            for path in sorted(root.rglob("*")):
                if not path.is_file():
                    continue
                if path.suffix.lower() not in {".docx", ".pptx", ".pdf", ".doc", ".txt", ".md", ".rtf"}:
                    continue
                try:
                    text = extract(path)
                    status = "ok" if text.strip() else "empty"
                except Exception as exc:
                    text = ""
                    status = f"error:{type(exc).__name__}"
                rows.append({"suite": suite, "file": str(path), "chars": str(len(text)), "status": status})
                if text.strip():
                    suite_parts.append(f"\n\n===== SOURCE FILE: {path} =====\n{text}")
        for name in F_OCR_MAP.get(suite, []):
            path = F_OCR_DIR / name
            if path.exists():
                text = path.read_text(encoding="utf-8", errors="ignore")
                rows.append({"suite": suite, "file": str(path), "chars": str(len(text)), "status": "prior_ocr"})
                suite_parts.append(f"\n\n===== PRIOR OCR FILE: {path} =====\n{text}")
        (OUT_DIR / f"{suite}.txt").write_text(clean("\n".join(suite_parts)), encoding="utf-8")
    with (RUN_DIR / "SOURCE_TEXT_INDEX.csv").open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=["suite", "file", "chars", "status"])
        writer.writeheader()
        writer.writerows(rows)
    summary = ["# SOURCE_TEXT_SUMMARY", ""]
    for suite in SUITE_DIRS:
        p = OUT_DIR / f"{suite}.txt"
        summary.append(f"- {suite}: {p.stat().st_size} bytes")
    (RUN_DIR / "SOURCE_TEXT_SUMMARY.md").write_text("\n".join(summary) + "\n", encoding="utf-8")


if __name__ == "__main__":
    main()

