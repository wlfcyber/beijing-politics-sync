#!/usr/bin/env python3
"""ClaudeCode independent extraction of 2026 二模 originals.

Reads .docx (python-docx), .pdf (PyMuPDF / fitz), .pptx (python-pptx),
and .doc (binary text scan). Writes one .txt per source under
02_claudecode_independent/cc_raw_text/<district>/<name>.txt
"""
from __future__ import annotations

import os
import sys
import re
from pathlib import Path

SRC_ROOT = Path(r"C:\Users\Administrator\Desktop\2026各区模拟题\2026各区二模")
OUT_ROOT = Path(
    r"C:\Users\Administrator\Desktop\02_代码项目与工具\mac-thread-restore\beijing-politics-sync-visible\reports\选必一_哲学宝典式重建_2026-05-16\08_2026_second_mock_backfill\02_claudecode_independent\cc_raw_text"
)


def extract_docx(p: Path) -> str:
    import docx
    d = docx.Document(str(p))
    parts = []
    for para in d.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in d.tables:
        for row in table.rows:
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    parts.append("[TBL] " + t)
    return "\n".join(parts)


def extract_pdf(p: Path) -> str:
    import fitz
    doc = fitz.open(str(p))
    parts = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        parts.append(f"--- page {i} ---\n{text}")
    doc.close()
    return "\n".join(parts)


TESSERACT_EXE = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
TESSDATA_DIR = r"C:\Users\Administrator\AppData\Local\tessdata"


def extract_pdf_ocr(p: Path, lang: str = "chi_sim+eng") -> str:
    """Fallback: render each page to image and run tesseract."""
    import fitz
    import subprocess
    import tempfile
    doc = fitz.open(str(p))
    parts = []
    tmpdir = tempfile.mkdtemp(prefix="cc_ocr_")
    for i, page in enumerate(doc, start=1):
        pix = page.get_pixmap(dpi=300)
        tmp_path = os.path.join(tmpdir, f"page_{i:03d}.png")
        pix.save(tmp_path)
        try:
            cp = subprocess.run(
                [TESSERACT_EXE, tmp_path, "-",
                 "--tessdata-dir", TESSDATA_DIR,
                 "-l", lang, "--psm", "6"],
                capture_output=True, text=True, encoding="utf-8",
            )
            parts.append(f"--- page {i} (OCR) ---\n{cp.stdout}")
            if cp.returncode != 0:
                parts.append(f"  [tesseract stderr: {cp.stderr[:200]}]")
        except FileNotFoundError:
            parts.append(f"--- page {i} (OCR FAIL: no tesseract) ---")
    doc.close()
    return "\n".join(parts)


def extract_pptx(p: Path) -> str:
    from pptx import Presentation
    pres = Presentation(str(p))
    parts = []
    for idx, slide in enumerate(pres.slides, start=1):
        parts.append(f"--- slide {idx} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs).strip()
                    if line:
                        parts.append(line)
            elif shape.shape_type == 19 and hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    parts.append(t)
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            parts.append("[TBL] " + t)
    return "\n".join(parts)


def extract_doc_binary(p: Path) -> str:
    """Very rough .doc text scan: pull printable runs of utf-16 LE.
    Most legacy .doc embed text as utf-16 LE in WordDocument stream.
    """
    raw = p.read_bytes()
    # Try utf-16 LE decoding by stripping high zero bytes
    try:
        s = raw.decode("utf-16-le", errors="ignore")
    except Exception:
        s = raw.decode("latin-1", errors="ignore")
    # Keep only CJK + ASCII printable, drop control characters
    keep = []
    buf = []
    for ch in s:
        cp = ord(ch)
        if (0x4E00 <= cp <= 0x9FFF) or (0x3400 <= cp <= 0x4DBF) or \
           (0x3000 <= cp <= 0x303F) or (0xFF00 <= cp <= 0xFFEF) or \
           (0x20 <= cp <= 0x7E) or ch in "\n\t":
            buf.append(ch)
        else:
            if buf:
                run = "".join(buf)
                if len(run.strip()) >= 2:
                    keep.append(run.strip())
                buf = []
    if buf:
        run = "".join(buf)
        if len(run.strip()) >= 2:
            keep.append(run.strip())
    text = "\n".join(keep)
    # Collapse very short fragments back together with spaces
    return text


def process_one(src: Path, out_dir: Path, log) -> None:
    ext = src.suffix.lower()
    out = out_dir / (src.stem + ext.replace(".", "_") + ".txt")
    out.parent.mkdir(parents=True, exist_ok=True)
    try:
        if ext == ".docx":
            text = extract_docx(src)
            method = "python-docx"
        elif ext == ".pdf":
            text = extract_pdf(src)
            method = "pymupdf text"
            # If essentially empty -> OCR
            cleaned = re.sub(r"\s|--- page \d+ ---", "", text)
            if len(cleaned) < 100:
                log.append(
                    f"  ! {src.name}: text layer too thin ({len(cleaned)} cjk-ish chars after strip), retrying with OCR"
                )
                text2 = extract_pdf_ocr(src)
                if text2 and "OCR FAIL" not in text2:
                    text = text2
                    method = "pymupdf -> tesseract OCR"
        elif ext == ".pptx":
            text = extract_pptx(src)
            method = "python-pptx"
        elif ext == ".doc":
            text = extract_doc_binary(src)
            method = "binary utf-16 scan"
        else:
            log.append(f"  ? {src.name}: unsupported ext {ext}")
            return
        out.write_text(text, encoding="utf-8")
        log.append(
            f"  OK {src.name} -> {out.name} ({method}; {len(text)} chars)"
        )
    except Exception as e:
        log.append(f"  FAIL {src.name} ({ext}): {e!r}")


def main() -> int:
    OUT_ROOT.mkdir(parents=True, exist_ok=True)
    log: list[str] = []
    log.append(f"src_root={SRC_ROOT}")
    log.append(f"out_root={OUT_ROOT}")
    for district_dir in sorted(SRC_ROOT.iterdir()):
        if not district_dir.is_dir():
            continue
        district = district_dir.name
        log.append(f"## {district}")
        out_dir = OUT_ROOT / district
        for root, dirs, files in os.walk(district_dir):
            for f in sorted(files):
                p = Path(root) / f
                # mirror sub-directory under out_dir (试卷/细则)
                rel = Path(root).relative_to(district_dir)
                sub_out = out_dir / rel
                process_one(p, sub_out, log)
    log_path = OUT_ROOT / "_extraction_log.txt"
    log_path.write_text("\n".join(log), encoding="utf-8")
    print("\n".join(log))
    return 0


if __name__ == "__main__":
    sys.exit(main())
