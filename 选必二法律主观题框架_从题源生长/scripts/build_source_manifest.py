#!/usr/bin/env python3
import csv
import hashlib
import json
import os
import re
import subprocess
import sys
import traceback
import zipfile
from collections import defaultdict
from pathlib import Path
from xml.etree import ElementTree as ET

RUN_ROOT = Path("/Users/wanglifei/Desktop/北京高考政治/选必二法律主观题框架_从题源生长")
MANIFEST_DIR = RUN_ROOT / "00_manifest"
TEXT_DIR = MANIFEST_DIR / "extracted_text"
RENDER_DIR = MANIFEST_DIR / "rendered_pages"
FAILED_PATH = MANIFEST_DIR / "failed_files.csv"
MANIFEST_CSV = MANIFEST_DIR / "source_manifest.csv"
MANIFEST_JSONL = MANIFEST_DIR / "source_manifest.jsonl"
PROCESSING_LOG = MANIFEST_DIR / "processing_log.md"

SOURCE_ROOTS = [
    Path("/Users/wanglifei/Desktop/2024模拟题"),
    Path("/Users/wanglifei/Desktop/2025模拟题"),
    Path("/Users/wanglifei/Desktop/2026模拟题"),
    Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中"),
    Path("/Users/wanglifei/GaokaoPolitics/2024各区模拟题"),
    Path("/Users/wanglifei/GaokaoPolitics/2025各区模拟题"),
    Path("/Users/wanglifei/GaokaoPolitics/2026各区模拟题"),
]

RAW_ARCHIVE_EXTS = {".zip", ".rar", ".7z"}
SUPPORTED_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".rtf", ".xlsx"}
SKIP_NAMES = {".DS_Store"}

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def clean_text(text: str) -> str:
    text = text.replace("\u0000", "")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()


def sha256_path(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def infer_year(path: Path) -> str:
    m = re.search(r"(20[2-3][0-9])", str(path))
    return m.group(1) if m else ""


def infer_district(path: Path) -> str:
    text = str(path)
    districts = [
        "东城", "西城", "朝阳", "海淀", "丰台", "石景山", "顺义", "通州", "昌平",
        "房山", "门头沟", "延庆", "大兴", "密云", "怀柔", "平谷", "燕山",
    ]
    for district in districts:
        if district in text:
            return district
    return ""


def infer_exam_stage(path: Path) -> str:
    text = str(path)
    patterns = [
        ("一模", "一模"),
        ("二模", "二模"),
        ("三模", "三模"),
        ("期中", "期中"),
        ("期末", "期末"),
        ("适应", "适应性练习"),
        ("综合练习（一）", "一模"),
        ("综合练习(一)", "一模"),
        ("综合练习（二）", "二模"),
        ("综合练习(二)", "二模"),
    ]
    for key, value in patterns:
        if key in text:
            return value
    return ""


def infer_role(path: Path) -> str:
    text = str(path).lower()
    name = path.name.lower()
    if "~$" in name:
        return "unknown"
    if "分类" in text or "汇编" in text:
        return "user_framework"
    if any(k in text for k in ["阅卷总结", "阅卷报告", "总结"]):
        return "marking_report"
    if any(k in text for k in ["评标", "评价标准"]):
        return "evaluation_standard"
    if any(k in text for k in ["讲评", "课件"]):
        return "lecture"
    if any(k in text for k in ["细则", "评分", "阅卷细则", "分题细则"]):
        return "subjective_rubric"
    if any(k in text for k in ["答案表", "客观题", "选择题答案"]):
        return "objective_answer_sheet"
    if any(k in text for k in ["答案", "参考"]):
        return "answer"
    if any(k in text for k in ["/试卷/", "试卷", "试题"]):
        return "paper"
    return "unknown"


def suite_key(path: Path) -> str:
    year = infer_year(path)
    district = infer_district(path)
    stage = infer_exam_stage(path)
    return "|".join([year, district, stage])


def extract_docx(path: Path):
    if Document is None:
        return "", 0, "python-docx_unavailable"
    doc = Document(str(path))
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return clean_text("\n".join(parts)), len(doc.paragraphs), "python-docx"


def extract_doc_or_rtf(path: Path):
    proc = subprocess.run(
        ["/usr/bin/textutil", "-convert", "txt", "-stdout", str(path)],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=45,
    )
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.strip() or "textutil failed")
    return clean_text(proc.stdout), "", "textutil"


def extract_pptx(path: Path):
    if Presentation is None:
        return "", 0, "python-pptx_unavailable"
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"[slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return clean_text("\n".join(parts)), len(prs.slides), "python-pptx"


def extract_xlsx_xml(path: Path):
    # Minimal xlsx text extraction without relying on openpyxl.
    shared = []
    parts = []
    with zipfile.ZipFile(path) as zf:
        if "xl/sharedStrings.xml" in zf.namelist():
            root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
            ns = {"a": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
            for si in root.findall("a:si", ns):
                texts = [t.text or "" for t in si.findall(".//a:t", ns)]
                shared.append("".join(texts))
        for name in sorted(n for n in zf.namelist() if n.startswith("xl/worksheets/sheet") and n.endswith(".xml")):
            root = ET.fromstring(zf.read(name))
            ns = {"a": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
            parts.append(f"[{name}]")
            for row in root.findall(".//a:row", ns):
                cells = []
                for c in row.findall("a:c", ns):
                    t = c.attrib.get("t")
                    v = c.find("a:v", ns)
                    value = ""
                    if v is not None and v.text is not None:
                        if t == "s":
                            idx = int(v.text)
                            value = shared[idx] if idx < len(shared) else v.text
                        else:
                            value = v.text
                    if value:
                        cells.append(value)
                if cells:
                    parts.append(" | ".join(cells))
    return clean_text("\n".join(parts)), "", "xlsx-xml"


def extract_pdf(path: Path, file_id: str):
    if fitz is None:
        return "", "", "pdf_unavailable", "unknown"
    doc = fitz.open(str(path))
    page_count = doc.page_count
    parts = []
    image_blocks = 0
    for i in range(page_count):
        page = doc.load_page(i)
        text = page.get_text("text") or ""
        if text.strip():
            parts.append(f"[page {i + 1}]\n{text.strip()}")
        try:
            image_blocks += len(page.get_images(full=True))
        except Exception:
            pass
    text = clean_text("\n\n".join(parts))
    if text:
        status = "text_layer"
    else:
        status = "image_pdf" if image_blocks else "no_text_layer"
        out_dir = RENDER_DIR / file_id
        out_dir.mkdir(parents=True, exist_ok=True)
        for i in range(min(page_count, 30)):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(1.4, 1.4), alpha=False)
            pix.save(str(out_dir / f"page_{i + 1:03d}.png"))
        if page_count > 30:
            (out_dir / "RENDER_TRUNCATED_AFTER_30_PAGES.txt").write_text(
                f"Rendered first 30 pages out of {page_count}; OCR tool unavailable.\n",
                encoding="utf-8",
            )
    doc.close()
    return text, page_count, "pymupdf_text_extract", status


def write_text_file(file_id: str, text: str) -> str:
    out = TEXT_DIR / f"{file_id}.txt"
    out.write_text(text + ("\n" if text else ""), encoding="utf-8")
    return str(out)


def main():
    MANIFEST_DIR.mkdir(parents=True, exist_ok=True)
    TEXT_DIR.mkdir(parents=True, exist_ok=True)
    RENDER_DIR.mkdir(parents=True, exist_ok=True)

    files = []
    seen = set()
    missing_roots = []
    for root in SOURCE_ROOTS:
        if not root.exists():
            missing_roots.append(str(root))
            continue
        for p in root.rglob("*"):
            if not p.is_file():
                continue
            if str(p) in seen:
                continue
            seen.add(str(p))
            files.append(p)

    rows = []
    failed = []
    suite_to_ids = defaultdict(list)
    hash_to_ids = defaultdict(list)
    archive_notes = []

    for idx, path in enumerate(sorted(files, key=lambda x: str(x)), start=1):
        file_id = f"F{idx:04d}"
        ext = path.suffix.lower()
        role = infer_role(path)
        text_status = "unknown"
        page_count = ""
        processing_method = "not_processed"
        notes = []
        working_path = str(path)
        text_out = ""
        digest = ""

        try:
            digest = sha256_path(path)
        except Exception as e:
            digest = ""
            notes.append(f"hash_failed: {e}")

        try:
            if path.name in SKIP_NAMES:
                text_status = "unknown"
                notes.append("macOS metadata file skipped")
            elif ext in RAW_ARCHIVE_EXTS:
                archive_notes.append(str(path))
                text_status = "unknown"
                processing_method = "archive_not_unpacked_in_manifest_step"
                notes.append("archive found; no raw-source archive unpacked unless evidence-root approved")
            elif ext == ".pdf":
                text, page_count, processing_method, text_status = extract_pdf(path, file_id)
                text_out = write_text_file(file_id, text)
                if text_status in {"image_pdf", "no_text_layer"}:
                    failed.append({
                        "file_id": file_id,
                        "file_path": str(path),
                        "failure_stage": "ocr",
                        "failure_reason": "PDF has no extractable text layer; pages rendered but local OCR unavailable",
                        "suggested_fix": "Run OCR/vision on rendered_pages for this file before evidence extraction",
                    })
            elif ext == ".docx":
                text, page_count, processing_method = extract_docx(path)
                text_status = "text_layer" if text else "no_text_layer"
                text_out = write_text_file(file_id, text)
            elif ext in {".doc", ".rtf"}:
                text, page_count, processing_method = extract_doc_or_rtf(path)
                text_status = "old_word" if ext == ".doc" else "text_layer"
                text_out = write_text_file(file_id, text)
            elif ext == ".pptx":
                text, page_count, processing_method = extract_pptx(path)
                text_status = "ppt"
                text_out = write_text_file(file_id, text)
            elif ext == ".xlsx":
                text, page_count, processing_method = extract_xlsx_xml(path)
                text_status = "text_layer" if text else "unknown"
                text_out = write_text_file(file_id, text)
            else:
                text_status = "unknown"
                notes.append(f"unsupported extension: {ext}")
                if ext not in {"", ".ds_store"}:
                    failed.append({
                        "file_id": file_id,
                        "file_path": str(path),
                        "failure_stage": "extract",
                        "failure_reason": f"Unsupported extension {ext}",
                        "suggested_fix": "Inspect manually if this file is a source artifact",
                    })
        except Exception as e:
            text_status = "corrupted"
            processing_method = processing_method + "_failed"
            notes.append(f"extract_failed: {e}")
            failed.append({
                "file_id": file_id,
                "file_path": str(path),
                "failure_stage": "extract",
                "failure_reason": repr(e),
                "suggested_fix": "Open/render manually or convert with Office/OCR",
            })

        skey = suite_key(path)
        suite_to_ids[skey].append(file_id)
        if digest:
            hash_to_ids[digest].append(file_id)

        row = {
            "file_id": file_id,
            "original_file_path": str(path),
            "working_file_path": working_path,
            "file_name": path.name,
            "year": infer_year(path),
            "district": infer_district(path),
            "exam_stage": infer_exam_stage(path),
            "file_type": ext.lstrip(".") or "unknown",
            "suspected_role": role,
            "text_layer_status": text_status,
            "page_count": page_count,
            "hash": digest,
            "processing_method": processing_method,
            "linked_files": "",
            "notes": "; ".join(notes + ([f"extracted_text={text_out}"] if text_out else [])),
        }
        rows.append(row)

    # Linked files: same year/district/stage and exact duplicates.
    id_to_links = defaultdict(set)
    for ids in suite_to_ids.values():
        if len(ids) > 1:
            for fid in ids:
                id_to_links[fid].update(x for x in ids if x != fid)
    for ids in hash_to_ids.values():
        if len(ids) > 1:
            for fid in ids:
                id_to_links[fid].update(x for x in ids if x != fid)
    for row in rows:
        row["linked_files"] = "|".join(sorted(id_to_links[row["file_id"]]))

    fieldnames = [
        "file_id",
        "original_file_path",
        "working_file_path",
        "file_name",
        "year",
        "district",
        "exam_stage",
        "file_type",
        "suspected_role",
        "text_layer_status",
        "page_count",
        "hash",
        "processing_method",
        "linked_files",
        "notes",
    ]

    with MANIFEST_CSV.open("w", newline="", encoding="utf-8-sig") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)

    with MANIFEST_JSONL.open("w", encoding="utf-8") as f:
        for row in rows:
            f.write(json.dumps(row, ensure_ascii=False) + "\n")

    with FAILED_PATH.open("w", newline="", encoding="utf-8-sig") as f:
        fieldnames_failed = ["file_id", "file_path", "failure_stage", "failure_reason", "suggested_fix"]
        writer = csv.DictWriter(f, fieldnames=fieldnames_failed)
        writer.writeheader()
        writer.writerows(failed)

    role_counts = defaultdict(int)
    status_counts = defaultdict(int)
    ext_counts = defaultdict(int)
    for row in rows:
        role_counts[row["suspected_role"]] += 1
        status_counts[row["text_layer_status"]] += 1
        ext_counts[row["file_type"]] += 1

    PROCESSING_LOG.write_text(
        "\n".join([
            "# Source Manifest Processing Log",
            "",
            f"generated_at: 2026-05-19T00:00:00+08:00",
            f"source_files_seen: {len(rows)}",
            f"failed_or_ocr_gap_rows: {len(failed)}",
            "",
            "## Source Roots",
            *[f"- {root} :: {'exists' if root.exists() else 'missing'}" for root in SOURCE_ROOTS],
            "",
            "## Missing Roots",
            *([f"- {root}" for root in missing_roots] or ["- none"]),
            "",
            "## File Type Counts",
            *[f"- {k}: {v}" for k, v in sorted(ext_counts.items())],
            "",
            "## Suspected Role Counts",
            *[f"- {k}: {v}" for k, v in sorted(role_counts.items())],
            "",
            "## Text Layer Status Counts",
            *[f"- {k}: {v}" for k, v in sorted(status_counts.items())],
            "",
            "## Archive Notes",
            *([f"- {p}" for p in archive_notes] or ["- no archives found inside approved raw source roots"]),
            "",
            "## Tool Notes",
            "- PDF extraction: PyMuPDF text extraction.",
            "- DOC/DOCX/RTF/PPTX extraction: python-docx, python-pptx, or macOS textutil.",
            "- OCR: tesseract is unavailable in PATH; text-empty PDFs were rendered to 00_manifest/rendered_pages and recorded in failed_files.csv.",
            "- Old 选必二 generated artifacts are not used as evidence in this manifest step.",
            "",
        ]),
        encoding="utf-8",
    )

    print(json.dumps({
        "rows": len(rows),
        "failed_or_ocr_gap_rows": len(failed),
        "manifest_csv": str(MANIFEST_CSV),
        "manifest_jsonl": str(MANIFEST_JSONL),
        "failed_files": str(FAILED_PATH),
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    try:
        main()
    except Exception:
        traceback.print_exc()
        sys.exit(1)

