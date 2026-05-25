# -*- coding: utf-8 -*-
from __future__ import annotations

import csv
import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import fitz  # PyMuPDF


RUN = Path(__file__).resolve().parents[1]
INV = RUN / "01_source_inventory"
BUNDLE_DIR = INV / "suite_source_bundles"
BUNDLE_DIR.mkdir(parents=True, exist_ok=True)


def safe_name(s: str) -> str:
    return re.sub(r"[^\w\u4e00-\u9fff.-]+", "_", s).strip("_")


def docx_text(path: Path) -> str:
    try:
        from docx import Document

        doc = Document(str(path))
        parts: list[str] = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception:
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        chunks: list[str] = []
        with zipfile.ZipFile(path) as zf:
            root = ET.fromstring(zf.read("word/document.xml"))
            for para in root.findall(".//w:p", ns):
                text = "".join(t.text or "" for t in para.findall(".//w:t", ns))
                if text:
                    chunks.append(text)
        return "\n".join(chunks)


def pptx_text(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            parts.append(f"[slide {idx}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def pdf_text(path: Path) -> str:
    doc = fitz.open(str(path))
    parts: list[str] = []
    for idx, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            parts.append(f"[page {idx}]\n{text}")
    return "\n\n".join(parts)


def old_doc_text(path: Path) -> str:
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0
    try:
        doc = word.Documents.Open(str(path), ReadOnly=True, AddToRecentFiles=False)
        try:
            return doc.Content.Text
        finally:
            doc.Close(False)
    finally:
        word.Quit()
        pythoncom.CoUninitialize()


def extract(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    try:
        if suffix == ".docx":
            return docx_text(path), "docx-text"
        if suffix == ".pptx":
            return pptx_text(path), "pptx-text"
        if suffix == ".pdf":
            return pdf_text(path), "pdf-text"
        if suffix == ".doc":
            return old_doc_text(path), "word-com-doc-text"
        return "", f"unsupported-{suffix}"
    except Exception as exc:
        return "", f"failed-{type(exc).__name__}: {exc}"


def main() -> int:
    inv_path = INV / "source_suite_inventory.csv"
    rows = list(csv.DictReader(inv_path.open("r", encoding="utf-8-sig", newline="")))
    ledger_rows: list[dict] = []
    for row in rows:
        suite = row["suite"]
        source_dir = Path(row["source_dir"])
        out = BUNDLE_DIR / f"{safe_name(suite)}.md"
        bundle: list[str] = [f"# {suite}", "", f"source_dir: {source_dir}", ""]
        for f in sorted([p for p in source_dir.rglob("*") if p.is_file()], key=lambda p: str(p)):
            text, status = extract(f)
            ledger_rows.append({
                "suite": suite,
                "path": str(f),
                "suffix": f.suffix.lower(),
                "status": status,
                "chars": len(text),
            })
            bundle.append(f"\n\n## FILE: {f.name}\n")
            bundle.append(f"path: {f}\nstatus: {status}\nchars: {len(text)}\n")
            if text.strip():
                bundle.append("```text\n" + text.strip()[:120000] + "\n```")
            else:
                bundle.append("NO_TEXT_EXTRACTED")
        out.write_text("\n".join(bundle), encoding="utf-8")

    with (INV / "source_text_extraction_ledger.csv").open("w", newline="", encoding="utf-8-sig") as fp:
        writer = csv.DictWriter(fp, fieldnames=["suite", "path", "suffix", "status", "chars"])
        writer.writeheader()
        writer.writerows(ledger_rows)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
