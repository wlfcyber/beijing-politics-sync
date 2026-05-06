#!/usr/bin/env python3
"""Extract text and render evidence for Phase 02 hard samples.

This script is intentionally evidence-only: it does not write student-facing
content. It creates raw text, render images, and a CSV trial matrix for the five
hard samples named in PHASE02_HARD_SAMPLE_PLAN.md.
"""

from __future__ import annotations

import csv
import json
import re
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


RUN = Path("/Users/wanglifei/Desktop/北京高考政治/选必三逻辑与思维_四线从0重跑_2026-05-04")
OUT = RUN / "02_extraction" / "hard_samples"
TEXT_OUT = OUT / "raw_text"
RENDER_OUT = OUT / "renders"
MATRIX = RUN / "05_coverage" / "phase02_hard_sample_trial_matrix.csv"
VISUAL_QUEUE = OUT / "visual_fallback_queue.md"


SAMPLES = [
    {
        "sample_id": "HS01",
        "suite": "2026顺义一模",
        "question": "Q19(2)",
        "module_target": "思维-科学思维",
        "sources": [
            ("paper", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区一模/2026顺义一模/试卷/试卷.pdf")),
            ("rubric", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区一模/2026顺义一模/细则/细则.pptx")),
        ],
        "patterns": ["19", "科学思维", "客观性", "预见性", "可检验"],
    },
    {
        "sample_id": "HS02",
        "suite": "2025海淀二模",
        "question": "Q20",
        "module_target": "思维-辩证思维",
        "sources": [
            ("paper", Path("/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025海淀二模/试卷/试卷.pdf")),
            ("rubric", Path("/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025海淀二模/细则/补充材料/2025年海淀二模评标实录.docx")),
        ],
        "patterns": ["20", "辩证", "分析", "综合", "整体", "动态", "否定", "质量互变"],
    },
    {
        "sample_id": "HS03",
        "suite": "2026朝阳期中",
        "question": "Q21(2)",
        "module_target": "思维-创新思维",
        "sources": [
            ("paper", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026朝阳期中/试卷/试卷.pdf")),
            ("rubric", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026朝阳期中/细则/细则.docx")),
            ("rubric_supplement", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026朝阳期中/细则/补充材料/细则.docx")),
        ],
        "patterns": ["21", "创新", "超前", "联想", "逆向", "发散", "聚合"],
    },
    {
        "sample_id": "HS04",
        "suite": "2026通州期末",
        "question": "Q11",
        "module_target": "思维选择题-抽象到具体",
        "sources": [
            ("paper", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026通州期末/试卷/试卷.pdf")),
            ("rubric", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026通州期末/细则/细则.pptx")),
        ],
        "patterns": ["11", "感性具体", "思维抽象", "思维具体", "选择"],
    },
    {
        "sample_id": "HS05",
        "suite": "2026东城期末",
        "question": "Q17(2)",
        "module_target": "推理/边界",
        "sources": [
            ("paper", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026东城期末/试卷/试卷.pdf")),
            ("rubric", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026东城期末/细则/细则.pptx")),
        ],
        "glob_sources": [
            ("rubric_split", Path("/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026东城期末/细则/分题细则"), "*.pptx")
        ],
        "patterns": ["17", "推理", "三段论", "假言", "形式逻辑", "必要条件", "充分条件"],
    },
]


def normalize(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text, encoding="utf-8")


def extract_pdf(path: Path, out_prefix: Path, patterns: list[str]) -> tuple[str, list[str]]:
    doc = fitz.open(path)
    parts: list[str] = []
    hits: list[str] = []
    render_dir = RENDER_OUT / out_prefix.name
    render_dir.mkdir(parents=True, exist_ok=True)
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        parts.append(f"\n\n===== PAGE {i} =====\n{text}")
        page_norm = normalize(text)
        if any(p in page_norm for p in patterns):
            hits.append(f"page {i}: {page_norm[:500]}")
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            pix.save(render_dir / f"page_{i:02d}.png")
    return "".join(parts), hits


def extract_docx(path: Path) -> str:
    doc = Document(str(path))
    chunks: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            chunks.append(para.text)
    for table_index, table in enumerate(doc.tables, start=1):
        chunks.append(f"\n[TABLE {table_index}]")
        for row in table.rows:
            chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(chunks)


def extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    chunks: list[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        chunks.append(f"\n===== SLIDE {slide_index} =====")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                chunks.append(shape.text)
            if getattr(shape, "has_table", False):
                chunks.append(f"[TABLE on slide {slide_index}]")
                for row in shape.table.rows:
                    chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(chunks)


def extract_any(path: Path, out_prefix: Path, patterns: list[str]) -> tuple[str, list[str], str]:
    if not path.exists():
        return "", [], "missing"
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            text, hits = extract_pdf(path, out_prefix, patterns)
            method = "fitz_text_plus_pattern_page_render"
        elif suffix == ".docx":
            text = extract_docx(path)
            hits = [normalize(line)[:500] for line in text.splitlines() if any(p in line for p in patterns)]
            method = "python_docx"
        elif suffix == ".pptx":
            text = extract_pptx(path)
            hits = [normalize(line)[:500] for line in text.splitlines() if any(p in line for p in patterns)]
            method = "python_pptx_text"
        else:
            text = ""
            hits = []
            method = f"unsupported_{suffix}"
    except Exception as exc:  # keep evidence workflow moving, but log the failure
        return f"EXTRACTION_ERROR: {type(exc).__name__}: {exc}", [], "error"
    write_text(TEXT_OUT / f"{out_prefix.name}.txt", text)
    return text, hits, method


def main() -> None:
    TEXT_OUT.mkdir(parents=True, exist_ok=True)
    RENDER_OUT.mkdir(parents=True, exist_ok=True)
    rows: list[dict[str, str]] = []
    visual_lines = ["# Phase 02 Visual Fallback Queue", ""]
    for sample in SAMPLES:
        all_sources = list(sample["sources"])
        for kind, folder, glob_pattern in sample.get("glob_sources", []):
            if folder.exists():
                for path in sorted(folder.glob(glob_pattern)):
                    all_sources.append((kind, path))

        for source_role, path in all_sources:
            source_id = f"{sample['sample_id']}_{source_role}_{len(rows)+1:03d}"
            out_prefix = Path(f"{sample['sample_id']}__{source_role}__{path.stem}")
            text, hits, method = extract_any(path, out_prefix, sample["patterns"])
            char_count = len(text)
            hit_count = len(hits)
            status = "missing" if method == "missing" else ("needs_visual_review" if hit_count == 0 or char_count < 200 else "text_extracted")
            hit_file = TEXT_OUT / f"{out_prefix.name}.hits.json"
            write_text(hit_file, json.dumps(hits, ensure_ascii=False, indent=2))
            rows.append({
                "sample_id": sample["sample_id"],
                "suite": sample["suite"],
                "question": sample["question"],
                "module_target": sample["module_target"],
                "source_role": source_role,
                "source_path": str(path),
                "method": method,
                "char_count": str(char_count),
                "pattern_hit_count": str(hit_count),
                "status": status,
                "text_file": str(TEXT_OUT / f"{out_prefix.name}.txt"),
                "hit_file": str(hit_file),
            })
            if status != "text_extracted":
                visual_lines.append(f"- {sample['sample_id']} {source_role}: `{path}` -> {status}, method={method}, chars={char_count}, hits={hit_count}")

    MATRIX.parent.mkdir(parents=True, exist_ok=True)
    with MATRIX.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        writer.writerows(rows)
    write_text(VISUAL_QUEUE, "\n".join(visual_lines) + "\n")
    print(f"wrote {MATRIX}")
    print(f"rows={len(rows)}")
    print(f"visual_queue={VISUAL_QUEUE}")


if __name__ == "__main__":
    main()
