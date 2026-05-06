#!/usr/bin/env python3
from __future__ import annotations

import csv
import glob
import json
import re
import subprocess
from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation


ROOT = Path("/Users/wanglifei/Desktop/北京高考政治/选必三逻辑与思维_四线从0重跑_2026-05-04")
QUEUE = ROOT / "01_source_inventory" / "PRIORITY_SOURCE_QUEUE.md"
OUT = ROOT / "02_extraction" / "priority_queue_sources"
TEXT_DIR = OUT / "text"
RENDER_DIR = OUT / "renders"

PATTERNS = [
    "逻辑与思维",
    "科学思维",
    "辩证思维",
    "创新思维",
    "超前思维",
    "联想思维",
    "发散",
    "聚合",
    "逆向",
    "分析与综合",
    "质量互变",
    "辩证否定",
    "感性具体",
    "思维抽象",
    "思维具体",
    "三段论",
    "假言",
    "选言",
    "联言",
    "矛盾律",
    "排中律",
    "同一律",
    "换质",
    "换位",
    "类比",
    "归纳",
    "推理",
    "判断",
    "概念",
]


def slug(path: Path, idx: int) -> str:
    raw = str(path).replace("/Users/wanglifei/", "").replace("/", "_")
    raw = re.sub(r"[^\w\u4e00-\u9fff.-]+", "_", raw)
    return f"{idx:03d}_{raw[:180]}"


def queue_paths() -> list[Path]:
    text = QUEUE.read_text(encoding="utf-8")
    raw_items = re.findall(r"`([^`]+)`", text)
    paths: list[Path] = []
    for item in raw_items:
        # Skip comments or repeated references that are not absolute paths.
        if not item.startswith("/Users/wanglifei/"):
            continue
        if "*" in item:
            paths.extend(Path(p) for p in sorted(glob.glob(item)))
        else:
            paths.append(Path(item))
    seen = set()
    unique: list[Path] = []
    for p in paths:
        key = str(p)
        if key not in seen:
            seen.add(key)
            unique.append(p)
    return unique


def hit_rows(text: str) -> list[dict]:
    rows = []
    lines = text.splitlines()
    for i, line in enumerate(lines, start=1):
        matched = [pat for pat in PATTERNS if pat in line]
        if matched:
            rows.append({"line": i, "patterns": matched, "text": line[:500]})
    return rows


def extract_pdf(path: Path, sid: str) -> tuple[str, dict]:
    doc = fitz.open(path)
    chunks = []
    hit_pages: set[int] = set()
    for page_no, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        chunks.append(f"===== PAGE {page_no} =====\n{text.rstrip()}\n")
        if any(pat in text for pat in PATTERNS):
            hit_pages.add(page_no)
    full_text = "\n".join(chunks)
    render_dir = RENDER_DIR / sid
    render_dir.mkdir(parents=True, exist_ok=True)
    pages_to_render = sorted(hit_pages)
    if len(full_text.strip()) < 500 and len(doc) <= 14:
        pages_to_render = list(range(1, len(doc) + 1))
    for page_no in pages_to_render[:30]:
        page = doc[page_no - 1]
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
        pix.save(render_dir / f"page_{page_no:02d}.png")
    meta = {
        "method": "pymupdf_text_plus_render",
        "pages": len(doc),
        "rendered_pages": pages_to_render[:30],
        "weak_text_layer": len(full_text.strip()) < 500,
    }
    return full_text, meta


def extract_docx(path: Path) -> tuple[str, dict]:
    doc = Document(path)
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for t_i, table in enumerate(doc.tables, start=1):
        parts.append(f"===== TABLE {t_i} =====")
        for row in table.rows:
            parts.append(" | ".join(cell.text.replace("\n", " ").strip() for cell in row.cells))
    return "\n".join(parts), {"method": "python_docx", "paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}


def extract_pptx(path: Path) -> tuple[str, dict]:
    prs = Presentation(path)
    parts: list[str] = []
    for s_i, slide in enumerate(prs.slides, start=1):
        parts.append(f"===== SLIDE {s_i} =====")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
    return "\n".join(parts), {"method": "python_pptx", "slides": len(prs.slides)}


def extract_textutil(path: Path, sid: str) -> tuple[str, dict]:
    out_path = TEXT_DIR / f"{sid}.textutil.txt"
    cmd = ["textutil", "-convert", "txt", "-stdout", str(path)]
    proc = subprocess.run(cmd, text=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    if proc.returncode != 0:
        return "", {"method": "textutil", "error": proc.stderr.strip(), "returncode": proc.returncode}
    out_path.write_text(proc.stdout, encoding="utf-8")
    return proc.stdout, {"method": "textutil", "sidecar": str(out_path)}


def extract_one(path: Path, idx: int) -> dict:
    sid = slug(path, idx)
    row = {
        "source_id": sid,
        "path": str(path),
        "exists": path.exists(),
        "suffix": path.suffix.lower(),
        "status": "pending",
        "method": "",
        "char_count": 0,
        "hit_count": 0,
        "text_path": "",
        "hits_path": "",
        "render_dir": "",
        "error": "",
    }
    if not path.exists():
        row["status"] = "missing"
        return row
    try:
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            text, meta = extract_pdf(path, sid)
            row["render_dir"] = str(RENDER_DIR / sid)
        elif suffix == ".docx":
            text, meta = extract_docx(path)
        elif suffix == ".pptx":
            text, meta = extract_pptx(path)
        elif suffix in {".doc", ".rtf"}:
            text, meta = extract_textutil(path, sid)
        else:
            try:
                text = path.read_text(encoding="utf-8")
                meta = {"method": "plain_text"}
            except UnicodeDecodeError:
                text = ""
                meta = {"method": "unsupported_binary"}
        text_path = TEXT_DIR / f"{sid}.txt"
        hits_path = TEXT_DIR / f"{sid}.hits.json"
        text_path.write_text(text, encoding="utf-8")
        hits = hit_rows(text)
        hits_path.write_text(json.dumps(hits, ensure_ascii=False, indent=2), encoding="utf-8")
        row.update(
            {
                "status": "extracted" if text else "empty_or_failed",
                "method": meta.get("method", ""),
                "char_count": len(text),
                "hit_count": len(hits),
                "text_path": str(text_path),
                "hits_path": str(hits_path),
            }
        )
        if "error" in meta:
            row["error"] = meta["error"]
        return row
    except Exception as exc:
        row["status"] = "error"
        row["error"] = repr(exc)
        return row


def main() -> None:
    TEXT_DIR.mkdir(parents=True, exist_ok=True)
    RENDER_DIR.mkdir(parents=True, exist_ok=True)
    rows = [extract_one(path, idx) for idx, path in enumerate(queue_paths(), start=1)]
    manifest = OUT / "priority_queue_extraction_manifest.csv"
    with manifest.open("w", encoding="utf-8", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()) if rows else ["source_id"])
        writer.writeheader()
        writer.writerows(rows)
    print(f"wrote {manifest}")
    print(f"sources={len(rows)} extracted={sum(r['status']=='extracted' for r in rows)} missing={sum(r['status']=='missing' for r in rows)} errors={sum(r['status']=='error' for r in rows)}")


if __name__ == "__main__":
    main()
