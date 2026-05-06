#!/usr/bin/env python3
from __future__ import annotations

import json
import re
import subprocess
from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "02_extraction" / "priority_source_texts"

SOURCES = [
    # user framework already extracted separately
    "/Users/wanglifei/Desktop/2026模拟题/2026各区一模/2026顺义一模/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2026模拟题/2026各区一模/2026顺义一模/细则/细则.pptx",
    "/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026朝阳期中/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026朝阳期中/细则/细则.docx",
    "/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026朝阳期中/细则/补充材料/细则.docx",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025昌平二模/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025昌平二模/细则/细则.pptx",
    "/Users/wanglifei/Desktop/2024模拟题/2024朝阳期中/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2024模拟题/2024朝阳期中/试卷/补充材料/202411朝阳高三政治 期中1试题(1).pdf",
    "/Users/wanglifei/Desktop/2024模拟题/2024朝阳期中/细则/细则.docx",
    "/Users/wanglifei/Desktop/2024模拟题/2024朝阳期中/细则/补充材料/2024.11期中政治朝阳评标2.docx",
    "/Users/wanglifei/Desktop/2024模拟题/2024朝阳期中/细则/补充材料/细则.rtf",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区一模/2025朝阳一模/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区一模/2025朝阳一模/细则/补充材料/20250329高3阅卷总结17 1题 具身智能 任会波组 阐释论证.doc",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区一模/2025顺义一模/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区一模/2025顺义一模/细则/细则.docx",
    "/Users/wanglifei/GaokaoPolitics/2025各区模拟题/2025各区一模/2025顺义一模/2025北京顺义高三一模政治（教师版）.pdf",
    "/Users/wanglifei/GaokaoPolitics/2025各区模拟题/2025各区一模/2025顺义一模/顺义区2025届高三第一次模拟考试参考答案—评分细则.docx",
    "/Users/wanglifei/GaokaoPolitics/2025各区模拟题/2025各区二模/2025昌平二模/2025北京昌平高三二模政治（教师版）.pdf",
    "/Users/wanglifei/GaokaoPolitics/2025各区模拟题/2025各区二模/2025昌平二模/2025昌平二模 政治评标.pptx",
    "/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026通州期末/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026通州期末/细则/细则.pptx",
    "/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026海淀期末/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2026模拟题/2026各区期末和期中/2026海淀期末/细则/细则.pdf",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025海淀二模/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025海淀二模/细则/细则.docx",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025海淀二模/细则/补充材料/2025年海淀二模评标实录.docx",
    "/Users/wanglifei/Desktop/2025模拟题/2025各区二模/2025海淀二模/其他材料/2025届二模考试讲评0510.pdf",
    "/Users/wanglifei/Desktop/2024模拟题/2024东城一模/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2024模拟题/2024东城一模/细则/细则.pptx",
    "/Users/wanglifei/Desktop/2024模拟题/2024东城一模/细则/补充材料/北京市东城区2023-2024学年度第二学期高三综合练习（一）思想政治答案(1).pdf",
    "/Users/wanglifei/Desktop/2024模拟题/西城一模/试卷/试卷.docx",
    "/Users/wanglifei/Desktop/2024模拟题/西城一模/细则/补充材料/2024.4高三统一测试思想政治答案.docx",
    "/Users/wanglifei/Desktop/2024模拟题/海淀二模/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2024模拟题/海淀二模/细则/细则.docx",
    "/Users/wanglifei/Desktop/2024模拟题/海淀二模/细则/补充材料/高三二模：政治答案(2).docx",
    "/Users/wanglifei/Desktop/2024模拟题/2024朝阳二模/试卷/试卷.pdf",
    "/Users/wanglifei/Desktop/2024模拟题/2024朝阳二模/细则/细则.pdf",
]


def safe_name(path: Path) -> str:
    s = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff._-]+", "_", str(path))
    return s.strip("_")[-190:]


def extract_pdf(path: Path) -> str:
    doc = fitz.open(path)
    parts: list[str] = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text")
        parts.append(f"[[PAGE {i}]]\n{text}".rstrip())
    return "\n\n".join(parts)


def extract_docx(path: Path) -> str:
    doc = Document(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for ti, table in enumerate(doc.tables, 1):
        parts.append(f"[[TABLE {ti}]]")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " / ") for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    prs = Presentation(path)
    parts: list[str] = []
    for si, slide in enumerate(prs.slides, 1):
        parts.append(f"[[SLIDE {si}]]")
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    slide_text.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " / ") for cell in row.cells]
                    if any(cells):
                        slide_text.append(" | ".join(cells))
        parts.append("\n".join(slide_text))
    return "\n\n".join(parts)


def extract_textutil(path: Path) -> str:
    proc = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", str(path)],
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        check=False,
    )
    if proc.returncode != 0:
        raise RuntimeError(proc.stderr.strip())
    return proc.stdout


def extract(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix in {".rtf", ".doc"}:
        return extract_textutil(path)
    raise ValueError(f"unsupported: {path}")


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    manifest = []
    for raw in SOURCES:
        path = Path(raw)
        rec = {"source": raw, "exists": path.exists()}
        if not path.exists():
            rec["status"] = "missing"
            manifest.append(rec)
            continue
        out = OUT / f"{len(manifest)+1:03d}_{safe_name(path)}.txt"
        try:
            text = extract(path)
            header = f"source: {path}\nstatus: extracted\n\n"
            out.write_text(header + text, encoding="utf-8")
            rec.update({"status": "extracted", "output": str(out), "chars": len(text)})
        except Exception as e:
            rec.update({"status": "failed", "error": repr(e)})
        manifest.append(rec)
    (OUT / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    for rec in manifest:
        print(rec)


if __name__ == "__main__":
    main()
