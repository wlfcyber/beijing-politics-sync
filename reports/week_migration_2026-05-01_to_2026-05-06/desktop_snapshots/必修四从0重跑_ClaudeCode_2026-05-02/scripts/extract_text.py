#!/usr/bin/env python3
"""Extract plain text from PDF/DOCX/PPTX with stable layout.

Usage: extract_text.py <path> [--pages 1-3]
Reads source files only; never writes back to source paths.
"""
from __future__ import annotations

import argparse
import os
import sys
from pathlib import Path


def extract_pdf(path: Path, pages: str | None = None) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    n = doc.page_count
    rng = range(n)
    if pages:
        out: list[int] = []
        for chunk in pages.split(","):
            chunk = chunk.strip()
            if "-" in chunk:
                a, b = chunk.split("-", 1)
                out.extend(range(int(a) - 1, int(b)))
            else:
                out.append(int(chunk) - 1)
        rng = [i for i in out if 0 <= i < n]
    parts: list[str] = []
    for i in rng:
        try:
            page = doc.load_page(i)
            text = page.get_text("text") or ""
        except Exception as exc:  # noqa: BLE001
            text = f"[extract_error page={i + 1}: {exc}]"
        parts.append(f"--- page {i + 1} ---\n{text}")
    doc.close()
    return "\n\n".join(parts)


def render_pdf_to_images(path: Path, out_dir: Path, dpi: int = 200) -> list[Path]:
    import fitz

    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(path))
    written: list[Path] = []
    for i in range(doc.page_count):
        page = doc.load_page(i)
        pix = page.get_pixmap(dpi=dpi)
        out = out_dir / f"page_{i + 1:02d}.png"
        pix.save(str(out))
        written.append(out)
    doc.close()
    return written


def extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- slide {idx} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
            if shape.shape_type == 19 and shape.has_table:  # 19 == TABLE
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append("\t".join(cells))
    return "\n".join(parts)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("path")
    parser.add_argument("--pages", default=None, help="PDF only, e.g. 1-3,5")
    args = parser.parse_args()
    path = Path(args.path)
    if not path.exists():
        print(f"missing: {path}", file=sys.stderr)
        return 1
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        print(extract_pdf(path, args.pages))
    elif suffix == ".docx":
        print(extract_docx(path))
    elif suffix == ".pptx":
        print(extract_pptx(path))
    else:
        print(f"unsupported suffix: {suffix}", file=sys.stderr)
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
