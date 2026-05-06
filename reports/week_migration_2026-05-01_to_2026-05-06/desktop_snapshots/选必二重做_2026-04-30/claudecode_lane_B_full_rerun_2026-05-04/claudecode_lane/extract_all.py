#!/usr/bin/env python3
"""
Round 3 全量解析脚本 — 选必二《法律与生活》ClaudeCode 生产线 B
扫三个源目录，按文件类型用对应工具提取纯文本，输出 .txt + SOURCE_LEDGER.csv + EXTRACTION_LOG.md。
"""

import os, sys, csv, subprocess, traceback, re, hashlib, time
from pathlib import Path

PROJECT_ROOT = Path("/Users/wanglifei/Desktop/北京高考政治/选必二重做_2026-04-30/claudecode_lane_B_full_rerun_2026-05-04")
SOURCE_ROOTS = [
    Path("/Users/wanglifei/Desktop/2024模拟题"),
    Path("/Users/wanglifei/Desktop/2025模拟题"),
    Path("/Users/wanglifei/Desktop/2026模拟题"),
]
OUT_TXT_DIR = PROJECT_ROOT / "source_inventory" / "extracted"
LEDGER_CSV = PROJECT_ROOT / "00_control" / "SOURCE_LEDGER.csv"
LOG_MD = PROJECT_ROOT / "source_inventory" / "EXTRACTION_LOG.md"

OUT_TXT_DIR.mkdir(parents=True, exist_ok=True)

# 永久排除（per MASTER_REQUIREMENTS）
EXCLUDE_PATTERNS = [
    r"2026.*石景山.*期末",
]

L1_NAME_PAT = re.compile(r"评分细则|评标|阅卷总结|讲评|分题细则|答案变通|阅卷|监控|主观题细则|勿传")
L2_NAME_PAT = re.compile(r"答案|参考答案|评分参考")
L3_NAME_PAT = re.compile(r"选择题答案|客观题答案")
SRC_NAME_PAT = re.compile(r"试卷|原卷|原题")

def evidence_level(path: Path) -> str:
    s = str(path)
    if SRC_NAME_PAT.search(s) and not L1_NAME_PAT.search(s):
        return "SOURCE_PAPER"
    if L1_NAME_PAT.search(s):
        return "L1"
    if L3_NAME_PAT.search(s):
        return "L3"
    if L2_NAME_PAT.search(s):
        return "L2"
    return "UNKNOWN"

def is_excluded(path: Path) -> bool:
    s = str(path)
    return any(re.search(p, s) for p in EXCLUDE_PATTERNS)

def safe_relative(src: Path, root: Path) -> Path:
    try:
        return src.relative_to(root)
    except ValueError:
        return Path(src.name)

# --- extractors ---

def extract_textutil(src: Path) -> str:
    """For .doc / .docx / .rtf — uses macOS textutil"""
    res = subprocess.run(
        ["textutil", "-stdout", "-convert", "txt", str(src)],
        capture_output=True, timeout=120
    )
    if res.returncode != 0:
        raise RuntimeError(f"textutil failed: {res.stderr.decode('utf-8', errors='ignore')[:200]}")
    return res.stdout.decode("utf-8", errors="ignore")

def extract_docx(src: Path) -> str:
    """python-docx with table dedup (skip merged-cell duplicates)"""
    from docx import Document
    doc = Document(str(src))
    lines = []
    for p in doc.paragraphs:
        if p.text.strip():
            lines.append(p.text)
    for ti, t in enumerate(doc.tables):
        lines.append(f"\n--- TABLE {ti+1} ---")
        seen_rows = set()
        for row in t.rows:
            seen_cells = []
            seen_tcs = set()
            for c in row.cells:
                tc_id = id(c._tc)
                if tc_id in seen_tcs:
                    continue
                seen_tcs.add(tc_id)
                seen_cells.append(c.text.strip().replace("\n", " / "))
            row_str = " | ".join(seen_cells)
            if row_str and row_str not in seen_rows:
                seen_rows.add(row_str)
                lines.append(row_str)
    return "\n".join(lines)

def extract_pptx(src: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(src))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n===== Slide {i} =====")
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs)
                        if t.strip():
                            lines.append(t)
                if shape.has_table:
                    lines.append("--- TABLE ---")
                    for row in shape.table.rows:
                        cells = [c.text.strip().replace("\n", " / ") for c in row.cells]
                        lines.append(" | ".join(cells))
            except Exception:
                pass
    return "\n".join(lines)

def extract_pdf_digital(src: Path) -> tuple[str, int]:
    """Returns (text, page_count). May return mostly-empty if scanned."""
    import pdfplumber
    chunks = []
    with pdfplumber.open(str(src)) as pdf:
        n = len(pdf.pages)
        for i, page in enumerate(pdf.pages, 1):
            t = page.extract_text() or ""
            chunks.append(f"\n===== Page {i} =====\n{t}")
            try:
                tables = page.extract_tables()
                for ti, tbl in enumerate(tables, 1):
                    chunks.append(f"--- Page {i} Table {ti} ---")
                    for row in tbl:
                        chunks.append(" | ".join((c or "").replace("\n", " ") for c in row))
            except Exception:
                pass
    return "\n".join(chunks), n

_OCR = None
def get_ocr():
    global _OCR
    if _OCR is None:
        from rapidocr_onnxruntime import RapidOCR
        _OCR = RapidOCR()
    return _OCR

def extract_pdf_ocr(src: Path) -> str:
    import pdfplumber
    ocr = get_ocr()
    chunks = []
    with pdfplumber.open(str(src)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            try:
                img = page.to_image(resolution=200).original
                tmp = f"/tmp/_ocr_pg_{i}.png"
                img.save(tmp)
                result, _ = ocr(tmp)
                if result:
                    text = "\n".join(line[1] for line in result)
                    chunks.append(f"\n===== Page {i} (OCR) =====\n{text}")
                else:
                    chunks.append(f"\n===== Page {i} (OCR empty) =====")
                try:
                    os.remove(tmp)
                except Exception:
                    pass
            except Exception as e:
                chunks.append(f"\n===== Page {i} (OCR ERROR: {e}) =====")
    return "\n".join(chunks)

def extract_pdf(src: Path) -> tuple[str, dict]:
    """Try digital extract first; if total chars < 100/page average, fall back to OCR."""
    text, n = extract_pdf_digital(src)
    # heuristic: scanned if very few chars per page
    real_text_len = sum(len(line) for line in text.split("\n") if not line.startswith("===") and not line.startswith("---"))
    avg_per_page = real_text_len / max(n, 1)
    info = {"pages": n, "digital_chars": real_text_len, "avg_per_page": round(avg_per_page, 1)}
    if avg_per_page < 80:  # likely scanned
        info["ocr_used"] = True
        ocr_text = extract_pdf_ocr(src)
        return ocr_text, info
    info["ocr_used"] = False
    return text, info

# --- main ---

def gather_files() -> list[Path]:
    out = []
    skip_exts = {".DS_Store", ".traineddata"}
    for root in SOURCE_ROOTS:
        for p in root.rglob("*"):
            if not p.is_file():
                continue
            if p.name.startswith("."):
                continue
            if p.suffix.lower() in {".xlsx"}:
                # skip xlsx for now (rarely contains 选必二 question text)
                continue
            ext = p.suffix.lower()
            if ext in {".doc", ".docx", ".pptx", ".pdf", ".rtf"}:
                out.append(p)
    return out

def slug(s: str) -> str:
    s = s.replace("/", "_").replace(" ", "_")
    return re.sub(r"[^\w一-鿿().+-]", "_", s)[:120]

def main():
    files = gather_files()
    print(f"Discovered {len(files)} files across {len(SOURCE_ROOTS)} roots", flush=True)
    rows = []
    log_lines = ["# Extraction Log", f"\nStarted: {time.strftime('%Y-%m-%d %H:%M:%S')}", f"\nTotal files: {len(files)}\n"]
    excluded_count = 0
    success_count = 0
    fail_count = 0
    ocr_count = 0
    t0 = time.time()

    for idx, src in enumerate(files, 1):
        # determine year and suite from path
        try:
            for root in SOURCE_ROOTS:
                if str(src).startswith(str(root)):
                    rel = src.relative_to(root)
                    year = root.name.replace("模拟题", "")
                    parts = rel.parts
                    suite = parts[0] if len(parts) > 1 else "ROOT"
                    if len(parts) > 2 and parts[1] not in {"细则", "试卷", "其他材料", "答案", "补充材料", "分题细则", "评标", "讲评", "阅卷总结"}:
                        # nested suite layer
                        suite = parts[1] if "20" in parts[1] else suite
                    break
        except Exception:
            year, suite = "UNKNOWN", "UNKNOWN"

        ext = src.suffix.lower()
        ev_level = evidence_level(src)
        excluded = is_excluded(src)
        rel_str = str(src.relative_to(src.parents[len(src.parents)-len(src.relative_to(SOURCE_ROOTS[0] if str(src).startswith(str(SOURCE_ROOTS[0])) else SOURCE_ROOTS[1] if str(src).startswith(str(SOURCE_ROOTS[1])) else SOURCE_ROOTS[2]).parts)-1])) if False else str(src)

        out_subdir = OUT_TXT_DIR / year / slug(suite)
        out_subdir.mkdir(parents=True, exist_ok=True)
        # use full relative path slug to avoid name collision
        rel_to_root = None
        for root in SOURCE_ROOTS:
            if str(src).startswith(str(root)):
                rel_to_root = src.relative_to(root)
                break
        rel_slug = slug("__".join(rel_to_root.parts)) if rel_to_root else slug(src.name)
        out_file = out_subdir / (rel_slug + ".txt")

        status = "ok"
        err = ""
        char_count = 0
        pdf_info = {}
        ocr_used = False

        if excluded:
            status = "excluded"
            log_lines.append(f"\n[{idx}/{len(files)}] EXCLUDED  {src}")
            excluded_count += 1
        else:
            try:
                if ext == ".doc" or ext == ".rtf":
                    text = extract_textutil(src)
                elif ext == ".docx":
                    text = extract_textutil(src)
                    # if textutil yields short, fall back to python-docx
                    if len(text.strip()) < 50:
                        text = extract_docx(src)
                elif ext == ".pptx":
                    text = extract_pptx(src)
                elif ext == ".pdf":
                    text, pdf_info = extract_pdf(src)
                    ocr_used = pdf_info.get("ocr_used", False)
                    if ocr_used:
                        ocr_count += 1
                else:
                    text = ""
                    status = "skipped_ext"
                char_count = len(text)
                if text:
                    out_file.write_text(text, encoding="utf-8")
                    if status == "ok":
                        success_count += 1
                else:
                    status = "empty"
                    fail_count += 1
            except Exception as e:
                status = "error"
                err = f"{type(e).__name__}: {str(e)[:200]}"
                fail_count += 1
                log_lines.append(f"\n[{idx}/{len(files)}] ERROR  {src}\n    {err}")

        rows.append({
            "idx": idx,
            "year": year,
            "suite": suite,
            "ext": ext,
            "filename": src.name,
            "full_path": str(src),
            "out_txt": str(out_file) if status == "ok" else "",
            "char_count": char_count,
            "evidence_level_hypothesis": ev_level,
            "excluded": "Y" if excluded else "N",
            "ocr_used": "Y" if ocr_used else "N",
            "pdf_pages": pdf_info.get("pages", ""),
            "status": status,
            "error": err,
        })

        if idx % 10 == 0 or idx == len(files):
            elapsed = time.time() - t0
            print(f"  [{idx}/{len(files)}] elapsed={elapsed:.1f}s  ok={success_count} ocr={ocr_count} fail={fail_count} excl={excluded_count}", flush=True)

    # write CSV
    fieldnames = ["idx","year","suite","ext","filename","full_path","out_txt","char_count","evidence_level_hypothesis","excluded","ocr_used","pdf_pages","status","error"]
    with open(LEDGER_CSV, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=fieldnames)
        w.writeheader()
        w.writerows(rows)

    # summary
    log_lines.append(f"\n## Summary")
    log_lines.append(f"- Total files: {len(files)}")
    log_lines.append(f"- Success: {success_count}")
    log_lines.append(f"- Excluded: {excluded_count}")
    log_lines.append(f"- Failed: {fail_count}")
    log_lines.append(f"- OCR used: {ocr_count}")
    log_lines.append(f"- Wall clock: {time.time()-t0:.1f}s")

    LOG_MD.write_text("\n".join(log_lines), encoding="utf-8")
    print(f"\nDone. Ledger: {LEDGER_CSV}\nLog: {LOG_MD}")

if __name__ == "__main__":
    main()
