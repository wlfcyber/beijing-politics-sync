#!/usr/bin/env python3
"""Step 1: walk the three real source roots, extract text, build manifest."""
import os, sys, csv, hashlib, subprocess, traceback
from pathlib import Path

ROOTS = [
    Path("/Users/wanglifei/Desktop/2024模拟题"),
    Path("/Users/wanglifei/Desktop/2025模拟题"),
    Path("/Users/wanglifei/Desktop/2026模拟题"),
]
OUT_DIR = Path("/Users/wanglifei/Desktop/北京高考政治/claude_zero_run/extracted")
MANIFEST_CSV = Path("/Users/wanglifei/Desktop/北京高考政治/claude_zero_run/01_source_manifest.csv")
FAILED_CSV = Path("/Users/wanglifei/Desktop/北京高考政治/claude_zero_run/01_failed_files.csv")
LOG_MD = Path("/Users/wanglifei/Desktop/北京高考政治/claude_zero_run/01_processing_log.md")

EXCLUDE_DIRS = {"已放弃", "2026石景山期末"}

def short_hash(path: Path) -> str:
    return hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:12]

def role_of(path: Path) -> str:
    """Classify by parent dir first (more reliable than filename)."""
    parent_dirs = [p for p in path.parts[:-1]]
    stem = path.stem
    # First check: filename-level signals that override dir-level
    if "讲评" in stem or "评标" in stem or "阅卷" in stem or "教师版" in stem:
        return "rubric_like"
    if "答案" in stem and "试卷" not in stem and "试题" not in stem:
        return "rubric_like"
    if "细则" in stem:
        return "rubric"
    # Combined files
    if ("试卷" in stem or "试题" in stem) and "答案" in stem:
        return "paper_with_answer"
    # Dir-level
    if any(p == "细则" or p == "分题细则" for p in parent_dirs):
        return "rubric"
    if any(p == "试卷" for p in parent_dirs):
        # Could include 补充材料 subdir which often has 试题及答案 combined files
        if "补充材料" in parent_dirs and ("答案" in stem or "评分" in stem or "细则" in stem):
            return "paper_with_answer" if "试题" in stem or "试卷" in stem else "rubric_like"
        return "paper"
    if any(p == "其他材料" for p in parent_dirs):
        return "supplement"
    return "other"

def suite_of(path: Path) -> str:
    # Prefer the deepest path part matching a specific suite (year+region+stage)
    # Skip umbrella container names like "2025各区一模", "2026各区期末和期中"
    parts = path.parts
    # Detect the year from a year-only umbrella dir like "2024模拟题"
    year_hint = "?"
    for p in parts:
        for y in ("2024", "2025", "2026"):
            if p.startswith(y):
                year_hint = y
                break
    candidate = None
    for p in parts:
        if "各区" in p:
            candidate = candidate or p
            continue
        if any(k in p for k in ["一模", "二模", "期中", "期末", "思政"]) and (
            p.startswith("20") or "高三" in p or "顺义" in p or "丰台" in p or "东城" in p or "海淀" in p or "朝阳" in p or "西城" in p or "石景山" in p
        ):
            # If suite name doesn't have a year prefix, prepend it
            if not p.startswith(("2024","2025","2026")):
                return f"{year_hint}{p}"
            return p
    return candidate or "unknown_suite"

def year_of(suite: str) -> str:
    for y in ("2024", "2025", "2026"):
        if y in suite:
            return y
    return "?"

def extract_pdf(path: Path) -> str:
    import fitz
    doc = fitz.open(path)
    chunks = []
    for i, page in enumerate(doc, 1):
        try:
            t = page.get_text("text")
        except Exception:
            t = ""
        chunks.append(f"\n===== PAGE {i} =====\n{t}")
    return "".join(chunks).strip()

def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        parts.append(p.text)
    for t in doc.tables:
        for row in t.rows:
            parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts).strip()

def extract_doc(path: Path) -> str:
    # Use textutil to convert legacy .doc to .txt
    try:
        out = subprocess.run(["textutil", "-convert", "txt", "-stdout", str(path)],
                             capture_output=True, timeout=60)
        return out.stdout.decode("utf-8", errors="ignore").strip()
    except Exception as e:
        raise RuntimeError(f"textutil failed: {e}")

def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    p = Presentation(path)
    out = []
    for i, slide in enumerate(p.slides, 1):
        out.append(f"\n===== SLIDE {i} =====")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                out.append(shape.text)
    return "\n".join(out).strip()

def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    manifest_rows = []
    failed_rows = []
    log_lines = ["# 01 Processing Log", "", "## Source roots", ""]
    for r in ROOTS:
        log_lines.append(f"- {r}  (exists={r.exists()})")
    log_lines += ["", "## File-level results", ""]

    seen = 0
    by_ext = {}
    for root in ROOTS:
        for path in root.rglob("*"):
            if not path.is_file():
                continue
            # Exclude flagged dirs
            if any(part in EXCLUDE_DIRS for part in path.parts):
                continue
            ext = path.suffix.lower().lstrip(".")
            if ext not in {"pdf", "docx", "doc", "pptx", "ppt"}:
                continue
            by_ext[ext] = by_ext.get(ext, 0) + 1
            seen += 1
            sid = short_hash(path)
            suite = suite_of(path)
            year = year_of(suite)
            role = role_of(path)
            out_file = OUT_DIR / f"{sid}__{path.stem[:60]}.txt"
            status = "ok"
            chars = 0
            note = ""
            try:
                if ext == "pdf":
                    text = extract_pdf(path)
                elif ext == "docx":
                    text = extract_docx(path)
                elif ext == "doc":
                    text = extract_doc(path)
                elif ext == "pptx":
                    text = extract_pptx(path)
                elif ext == "ppt":
                    # convert via textutil best effort
                    text = extract_doc(path)
                else:
                    text = ""
                if not text.strip():
                    status = "empty_or_scan"
                    note = "no text layer; likely scan-only PDF (no OCR available)"
                out_file.write_text(text, encoding="utf-8")
                chars = len(text)
            except Exception as e:
                status = "fail"
                note = f"{type(e).__name__}: {e}"
                failed_rows.append([sid, str(path), ext, note])
            manifest_rows.append([
                sid, year, suite, role, ext, chars, status, str(path), out_file.name
            ])
            log_lines.append(f"- [{sid}] {ext} {chars} {status} {path.relative_to(root.parent) if path.is_relative_to(root.parent) else path}")

    with MANIFEST_CSV.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["source_id","year","suite","role","ext","chars","status","abs_path","extract_filename"])
        w.writerows(manifest_rows)

    with FAILED_CSV.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["source_id","abs_path","ext","error"])
        w.writerows(failed_rows)

    LOG_MD.write_text("\n".join(log_lines), encoding="utf-8")
    print(f"seen={seen} by_ext={by_ext}")
    print(f"failed={len(failed_rows)} manifest_rows={len(manifest_rows)}")

if __name__ == "__main__":
    main()
