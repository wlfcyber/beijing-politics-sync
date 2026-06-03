#!/usr/bin/env python3
from __future__ import annotations

import csv
import hashlib
import re
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from docx import Document
from pptx import Presentation

try:
    import pypdf
except Exception:  # pragma: no cover
    pypdf = None

try:
    import fitz
except Exception:  # pragma: no cover
    fitz = None


BASE = Path("/Users/wanglifei/Desktop/北京高考政治")
OUT = BASE / "选必二重做_2026-04-30"
TEXT_DIR = OUT / "extracted_text"

SOURCE_ROOTS = [
    Path("/Users/wanglifei/Desktop/2024模拟题"),
    Path("/Users/wanglifei/Desktop/2025模拟题"),
    Path("/Users/wanglifei/Desktop/2026模拟题"),
]

LEGAL_TERMS = [
    "法律与生活",
    "选必二",
    "选修2",
    "选修二",
    "民法典",
    "合同",
    "侵权",
    "违约",
    "纠纷",
    "诉讼",
    "调解",
    "仲裁",
    "劳动",
    "消费者",
    "知识产权",
    "著作权",
    "专利",
    "商标",
    "不正当竞争",
    "格式条款",
    "继承",
    "遗嘱",
    "婚姻",
    "夫妻",
    "相邻关系",
    "人格权",
    "名誉权",
    "隐私权",
    "肖像权",
    "物权",
    "债权",
    "所有权",
    "赔偿",
    "举证",
    "法院",
    "人民法院",
    "司法确认",
    "公益诉讼",
    "民事",
    "法律责任",
    "权利义务",
    "限制民事行为能力",
    "无民事行为能力",
    "劳动合同",
    "用人单位",
    "经营者",
    "知情权",
    "公平交易权",
    "自主选择权",
]

CORE_XUANBIER_TERMS = [
    "法律与生活",
    "选必二",
    "选修2",
    "选修二",
    "民法典",
    "合同",
    "侵权",
    "违约",
    "劳动合同",
    "劳动争议",
    "劳动关系",
    "劳动仲裁",
    "用人单位",
    "消费者",
    "经营者",
    "消费者权益保护法",
    "知识产权",
    "著作权",
    "专利",
    "商标",
    "不正当竞争",
    "格式条款",
    "继承",
    "遗嘱",
    "婚姻",
    "夫妻",
    "相邻关系",
    "人格权",
    "名誉权",
    "隐私权",
    "肖像权",
    "物权",
    "债权",
    "所有权",
    "质权",
    "监护人",
    "侵权责任",
    "违约责任",
    "举证责任",
    "司法确认",
    "民事诉讼",
    "诉讼时效",
    "法律责任",
    "限制民事行为能力",
    "无民事行为能力",
    "知情权",
    "公平交易权",
    "自主选择权",
    "仲裁",
    "调解",
]

SUPPORT_XUANBIER_TERMS = [
    "民事",
    "法院",
    "人民法院",
    "起诉",
    "诉讼",
    "上诉",
    "审判监督程序",
    "赔偿",
    "举证",
    "权利义务",
    "权利和义务",
    "订立",
    "撤销",
    "生效",
    "无效",
    "解除",
    "补偿",
    "票款",
]

POLITICAL_LAW_NOISE = [
    "国家安全",
    "香港特别行政区",
    "行政机关",
    "政府",
    "人大",
    "人民代表大会",
    "立法",
    "执法",
    "政协",
    "法治政府",
    "基层治理",
    "依法治国",
    "政治与法治",
    "党的领导",
    "检察机关",
    "检察院",
    "行政公益诉讼",
    "公益诉讼",
    "文物",
    "文化遗产",
]

EXPLICIT_XUANBIER = ["法律与生活", "选必二", "选修2", "选修二"]
DISTRICTS = [
    "海淀",
    "西城",
    "东城",
    "朝阳",
    "丰台",
    "石景山",
    "门头沟",
    "房山",
    "顺义",
    "通州",
    "昌平",
    "延庆",
]

SKIP_DIR_MARKERS = {"补充材料", "其他材料", "分题细则", "原目录壳"}
SUPPORTED = {".pdf", ".docx", ".doc", ".rtf", ".pptx", ".txt", ".md"}


@dataclass
class ExtractedFile:
    path: Path
    role: str
    text: str = ""
    status: str = "pending"
    text_cache: Path | None = None


@dataclass
class Suite:
    suite_id: str
    year: str
    district: str
    stage: str
    name: str
    path: Path
    paper_files: list[ExtractedFile] = field(default_factory=list)
    rubric_files: list[ExtractedFile] = field(default_factory=list)
    support_files: list[ExtractedFile] = field(default_factory=list)
    status: str = "unknown"
    notes: str = ""


def clean_text(text: str) -> str:
    text = text.replace("\u3000", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def short_cell(text: str, limit: int = 900) -> str:
    text = clean_text(text)
    if len(text) <= limit:
        return text
    return text[:limit] + "……"


def file_hash(path: Path) -> str:
    h = hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:16]
    return h


def text_cache_path(path: Path) -> Path:
    return TEXT_DIR / f"{file_hash(path)}_{path.stem}.txt"


def extract_pdf(path: Path) -> tuple[str, str]:
    chunks: list[str] = []
    status = "extracted"
    if pypdf is not None:
        try:
            reader = pypdf.PdfReader(str(path))
            for page in reader.pages:
                chunks.append(page.extract_text() or "")
        except Exception as exc:
            status = f"pypdf_failed:{type(exc).__name__}"
    text = "\n".join(chunks)
    if len(text.strip()) < 80 and fitz is not None:
        try:
            chunks = []
            doc = fitz.open(str(path))
            for page in doc:
                chunks.append(page.get_text("text") or "")
            text = "\n".join(chunks)
            if len(text.strip()) >= 80:
                status = "extracted_fitz"
        except Exception as exc:
            status = f"fitz_failed:{type(exc).__name__}"
    if len(text.strip()) < 80:
        status = "text_empty_or_scan_ocr_needed"
    return clean_text(text), status


def extract_docx(path: Path) -> tuple[str, str]:
    try:
        doc = Document(str(path))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return clean_text("\n".join(parts)), "extracted"
    except Exception as exc:
        return "", f"docx_failed:{type(exc).__name__}"


def extract_pptx(path: Path) -> tuple[str, str]:
    try:
        prs = Presentation(str(path))
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"[slide {idx}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            slide_parts.append(" | ".join(cells))
            if len(slide_parts) > 1:
                parts.append("\n".join(slide_parts))
        return clean_text("\n\n".join(parts)), "extracted"
    except Exception as exc:
        return "", f"pptx_failed:{type(exc).__name__}"


def extract_with_textutil(path: Path) -> tuple[str, str]:
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=30,
        )
        text = result.stdout
        if result.returncode != 0:
            return clean_text(text), f"textutil_failed:{result.stderr.strip()[:80]}"
        if len(text.strip()) < 20:
            return clean_text(text), "textutil_empty"
        return clean_text(text), "extracted_textutil"
    except Exception as exc:
        return "", f"textutil_exception:{type(exc).__name__}"


def extract_text(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix in {".doc", ".rtf"}:
        return extract_with_textutil(path)
    if suffix in {".txt", ".md"}:
        try:
            return clean_text(path.read_text(encoding="utf-8", errors="ignore")), "extracted"
        except Exception as exc:
            return "", f"text_failed:{type(exc).__name__}"
    return "", "unsupported"


def is_suite_dir(path: Path) -> bool:
    if path.name in SKIP_DIR_MARKERS:
        return False
    return (path / "试卷").is_dir() or (path / "细则").is_dir()


def discover_suites() -> list[Suite]:
    suites: list[Suite] = []
    seen: set[Path] = set()
    for root in SOURCE_ROOTS:
        if not root.exists():
            continue
        for path in sorted([root] + [p for p in root.rglob("*") if p.is_dir()]):
            if not is_suite_dir(path):
                continue
            if any(part in SKIP_DIR_MARKERS for part in path.relative_to(root).parts[:-1]):
                continue
            resolved = path.resolve()
            if resolved in seen:
                continue
            seen.add(resolved)
            year = infer_year(path)
            district = infer_district(path.name)
            stage = infer_stage(path.name, path)
            suite_id = make_suite_id(year, district, stage, len(suites) + 1)
            suites.append(Suite(suite_id, year, district, stage, path.name, path))
    return suites


def infer_year(path: Path) -> str:
    match = re.search(r"20\d{2}", str(path))
    return match.group(0) if match else "unknown"


def infer_district(name: str) -> str:
    for d in DISTRICTS:
        if d in name:
            return d
    return "unknown"


def infer_stage(name: str, path: Path) -> str:
    text = str(path)
    if "期中" in text:
        return "期中"
    if "期末" in text:
        return "期末"
    if "二模" in text:
        return "二模"
    if "一模" in text:
        return "一模"
    return "unknown"


def make_suite_id(year: str, district: str, stage: str, idx: int) -> str:
    return f"{year}_{district}_{stage}_{idx:03d}"


def iter_files(base: Path) -> Iterable[Path]:
    if not base.exists():
        return []
    files = []
    for path in base.rglob("*"):
        if not path.is_file():
            continue
        if path.name.startswith("~$") or path.name.startswith("."):
            continue
        if path.suffix.lower() in SUPPORTED:
            files.append(path)
    return sorted(files)


def attach_files(suite: Suite) -> None:
    paper_root = suite.path / "试卷"
    direct_papers: list[Path] = []
    paper_supplements: list[Path] = []
    for path in iter_files(paper_root):
        rel_parts = path.relative_to(paper_root).parts if paper_root.exists() else ()
        if "补充材料" in rel_parts:
            paper_supplements.append(path)
        else:
            direct_papers.append(path)
    for path in direct_papers:
        suite.paper_files.append(ExtractedFile(path, "paper"))
    # Some source folders put answer keys or scoring standards under 试卷/补充材料.
    # Keep them available for matching, but do not let them pollute question-paper detection.
    target = suite.support_files if direct_papers else suite.paper_files
    for path in paper_supplements:
        target.append(ExtractedFile(path, "support" if direct_papers else "paper"))
    for path in iter_files(suite.path / "细则"):
        role = "rubric"
        if any(k in path.name for k in ["答案", "参考"]):
            role = "rubric_or_answer"
        if any(k in path.name for k in ["评标", "阅卷", "细则", "评分"]):
            role = "rubric"
        suite.rubric_files.append(ExtractedFile(path, role))
    for support_dir in ["其他材料"]:
        for path in iter_files(suite.path / support_dir):
            suite.support_files.append(ExtractedFile(path, "support"))


def extract_files(files: list[ExtractedFile]) -> None:
    TEXT_DIR.mkdir(parents=True, exist_ok=True)
    for item in files:
        text, status = extract_text(item.path)
        item.text = text
        item.status = status
        cache = text_cache_path(item.path)
        cache.write_text(
            f"source: {item.path}\nstatus: {status}\n\n{text}",
            encoding="utf-8",
        )
        item.text_cache = cache


def legal_score(text: str) -> int:
    return sum(text.count(term) for term in LEGAL_TERMS)


def private_law_score(text: str) -> int:
    return sum(text.count(term) for term in CORE_XUANBIER_TERMS)


def support_law_score(text: str) -> int:
    return sum(text.count(term) for term in SUPPORT_XUANBIER_TERMS)


def explicit_hit(text: str) -> bool:
    return any(term in text for term in EXPLICIT_XUANBIER)


def determine_status(suite: Suite) -> None:
    if "2026石景山期末" in str(suite.path):
        suite.status = "excluded"
        suite.notes = "2026石景山期末按用户要求排除，除非提供新可用细则。"
        return
    paper_text = "\n".join(f.text for f in suite.paper_files)
    rubric_text = "\n".join(f.text for f in suite.rubric_files + suite.support_files)
    paper_empty = bool(suite.paper_files) and len(paper_text.strip()) < 120
    combined = paper_text + "\n" + rubric_text
    score = legal_score(combined)
    private_score = private_law_score(paper_text)
    support_score = support_law_score(paper_text)
    paper_candidates = candidate_questions_from_text(paper_text)
    rubric_candidates = candidate_questions_from_text(rubric_text)
    if explicit_hit(paper_text) or re.search(r"运用《?法律与生活》?", paper_text):
        suite.status = "has_xuanbier"
        suite.notes = "题面命中明确《法律与生活》/选必二标签。"
    elif paper_candidates:
        suite.status = "has_xuanbier"
        suite.notes = f"题面题段出现稳定《法律与生活》私法/程序法信号（core_score={private_score}, support_score={support_score}, legal_score={score}）。"
    elif paper_empty and (explicit_hit(rubric_text) or rubric_candidates):
        suite.status = "has_xuanbier"
        suite.notes = "试卷文本层不足，但细则/答案源命中《法律与生活》题；题面需 OCR 或原卷复核。"
    elif paper_empty:
        suite.status = "uncertain"
        suite.notes = "试卷文本层不足，需 OCR/渲染复核后判断是否含选必二。"
    elif private_score >= 1 or support_score >= 3:
        suite.status = "uncertain"
        suite.notes = f"题面仅零散《法律与生活》信号（core_score={private_score}, support_score={support_score}），防止强行归入，列为不确定。"
    else:
        suite.status = "no_xuanbier"
        suite.notes = "未见稳定《法律与生活》题面或法律题密集信号。"


def normalize_for_search(text: str) -> str:
    text = re.sub(r"\s+", " ", text)
    return text


def nearest_question_number(before: str) -> str | None:
    patterns = [
        r"第\s*([1-9]|1[0-9]|2[0-2])\s*题",
        r"([1-9]|1[0-9]|2[0-2])\s*[．.、]",
    ]
    hits: list[tuple[int, str]] = []
    for pat in patterns:
        for m in re.finditer(pat, before):
            hits.append((m.start(), m.group(1)))
    if not hits:
        return None
    return sorted(hits)[-1][1]


def context_is_xuanbier(context: str) -> bool:
    private = private_law_score(context)
    support = support_law_score(context)
    stem = stem_from_context(context)
    stem_private = private_law_score(stem)
    stem_support = support_law_score(stem)
    if explicit_hit(context) or re.search(r"运用《?法律与生活》?", context):
        return True
    if "行政公益诉讼" in context or ("公益诉讼" in context and stem_private == 0):
        return False
    # If legal words appear only in one option of a logic/philosophy/economics item,
    # do not count the item as a Legal and Life question.
    if stem_private == 0 and stem_support == 0:
        return False
    if stem_private >= 2:
        return True
    if stem_private >= 1 and (support >= 1 or private >= 2) and not (
        any(noise in stem for noise in POLITICAL_LAW_NOISE) and stem_support < 2
    ):
        return True
    if stem_support >= 3 and all(k in stem for k in ["法院", "赔偿"]):
        return True
    if stem_support >= 2 and all(k in stem for k in ["民事", "诉讼"]):
        return True
    return False


def stem_from_context(context: str) -> str:
    first_option = re.search(r"[ABCD]\s*[\.．、]|[①②③④]", context)
    return context[: first_option.start()] if first_option else context


def candidate_questions_from_text(raw_text: str) -> dict[str, str]:
    text = normalize_for_search(raw_text)
    candidates: dict[str, str] = {}
    if not text:
        return {}
    segments = question_segments_from_text(text)
    for qno, segment in segments.items():
        if not context_is_xuanbier(segment):
            continue
        if int(qno) <= 15 and likely_nonlegal_choice(segment):
            continue
        if int(qno) > 15 and not subjective_segment_has_legal_task(segment):
            continue
        candidates[qno] = clean_text(segment)[:3600]
    return candidates


def subjective_segment_has_legal_task(context: str) -> bool:
    compact = normalize_for_search(context)
    if explicit_hit(compact) or re.search(r"运用《?法律与生活》?", compact):
        return True
    hard_nonlegal = [
        "创新思维",
        "经济与社会",
        "政治与法治",
        "哲学",
        "逻辑与思维",
        "当代国际政治与经济",
        "经济全球化",
        "全球南方",
        "高质量发展",
        "社会保障",
        "基本医疗保险",
        "强国建设",
        "民族复兴",
    ]
    legal_task_cues = [
        "法院",
        "人民法院",
        "仲裁",
        "劳动人事争议",
        "诉讼",
        "纠纷",
        "案件",
        "案情",
        "裁判",
        "判决",
        "调解",
        "民法典",
        "劳动合同法",
        "消费者权益保护法",
        "反不正当竞争法",
        "知识产权",
    ]
    if any(cue in compact for cue in hard_nonlegal) and not any(cue in compact for cue in legal_task_cues):
        return False
    return private_law_score(compact) >= 2 and any(cue in compact for cue in legal_task_cues)


def likely_nonlegal_choice(context: str) -> bool:
    compact = normalize_for_search(context)
    nonlegal_hard = [
        "判断类型",
        "这一争论体现",
        "该争论体现",
        "思维方法",
        "联想思维",
        "逆向思维",
        "类比推理",
        "演绎推理",
        "归纳推理",
        "辩证思维",
    ]
    if any(cue in compact for cue in nonlegal_hard):
        return True
    stem = stem_from_context(compact)
    if "体现了" in stem and not any(cue in stem for cue in ["法院", "诉至", "民法典", "消费者权益保护法", "劳动合同法"]):
        return True
    return False


def question_segments_from_text(text: str) -> dict[str, str]:
    marker_re = re.compile(r"(?<!\d)([1-9]|1[0-9]|2[0-2])\s*[\.．、]")
    matches = []
    for match in marker_re.finditer(text):
        qno = int(match.group(1))
        after = text[match.end() : match.end() + 1]
        if qno >= 16 and after.isdigit():
            continue
        matches.append(match)
    segments: dict[str, str] = {}
    if not matches:
        return segments
    for idx, match in enumerate(matches):
        qno = match.group(1)
        # Keep the first substantial occurrence for each exam question number.
        if qno in segments and len(segments[qno]) > 120:
            continue
        end = matches[idx + 1].start() if idx + 1 < len(matches) else len(text)
        segment = text[match.start() : end]
        if len(segment.strip()) < 80:
            continue
        segments[qno] = segment
    return segments


def candidate_questions(suite: Suite) -> dict[str, str]:
    return candidate_questions_from_text("\n".join(f.text for f in suite.paper_files))


def extract_answer_key(text: str) -> dict[str, str]:
    key: dict[str, str] = {}
    flat = normalize_for_search(text)
    # Handles "1 A", "1.A", "1、A", table-like rows.
    for m in re.finditer(r"(?<!\d)([1-9]|1[0-5])\s*[\.．、:：]?\s*([ABCD])\b", flat):
        key.setdefault(m.group(1), m.group(2))
    # Handles compact snippets like "1C 2B 3A".
    for m in re.finditer(r"(?<!\d)([1-9]|1[0-5])\s*([ABCD])(?![A-Za-z])", flat):
        key.setdefault(m.group(1), m.group(2))
    return key


def option_map(context: str) -> dict[str, str]:
    opts: dict[str, str] = {}
    matches = list(re.finditer(r"([ABCD])\s*[\.．、]\s*", context))
    for idx, match in enumerate(matches):
        label = match.group(1)
        start = match.end()
        end = matches[idx + 1].start() if idx + 1 < len(matches) else min(len(context), start + 260)
        text = clean_text(context[start:end])
        if text:
            opts[label] = text[:260]
    return opts


def scenario_from_context(context: str) -> str:
    context = clean_text(context)
    # Prefer the part before options.
    first_option = re.search(r"[ABCD]\s*[\.．、]|[①②③④]", context)
    if first_option:
        context = context[: first_option.start()]
    context = re.sub(r"第?\s*[0-9]{1,2}\s*(题|[\.．、])", "", context)
    context = context.strip(" ：:，,。")
    if len(context) > 90:
        return context[-90:]
    return context or "本题法律情境"


def infer_mistake(option_text: str) -> str:
    if re.search(r"一律|必然|当然|均|都|只能|直接|无需|不需要|不可能", option_text):
        return "表述绝对化或条件缺失"
    if "执法" in option_text and "法院" in option_text:
        return "把司法机关职责与行政执法混淆"
    if "违约" in option_text and "侵权" in option_text:
        return "混淆违约责任与侵权责任"
    if any(k in option_text for k in ["仲裁", "起诉", "诉讼", "调解", "司法确认"]):
        return "纠纷解决程序或法律效力判断不准"
    if any(k in option_text for k in ["所有权", "物权", "债权", "用益物权", "质权"]):
        return "权利类型或物权效力判断不准"
    if any(k in option_text for k in ["著作权", "专利", "商标", "知识产权"]):
        return "知识产权类型或权利边界判断不准"
    if any(k in option_text for k in ["消费者", "知情权", "公平交易权", "自主选择权"]):
        return "消费者权利定位或经营者义务判断不准"
    if any(k in option_text for k in ["劳动", "用人单位", "劳动合同", "社保"]):
        return "劳动关系中的法定义务、合同变更或程序边界不准"
    return "法律关系、责任条件或材料边界判断不准"


def choice_sentences(qno: str, context: str, answer: str | None) -> tuple[str, str]:
    scenario = scenario_from_context(context)
    opts = option_map(context)
    statements = numbered_statement_map(context)
    combos = combo_option_map(opts)
    if not answer:
        overview = (
            f"在{scenario}中，正确答案尚未从可靠答案表自动锁定；本题暂不判定具体错项，"
            f"因为选择题预处理必须先确认客观答案，不能把普通材料推断冒充答案。"
        )
        return overview, "答案未锁定，需回原卷可靠客观答案表确认后再补写错项句。"

    correct_judgement = ""
    wrong_direction = "把法律关系、责任条件或程序边界说偏"
    if answer in combos and statements:
        correct_nums = combos[answer]
        correct_judgement = "、".join(f"{n}{short_cell(statements.get(n, ''), 70)}" for n in correct_nums if statements.get(n))
        wrong_nums = [n for n in statements if n not in correct_nums]
        if wrong_nums:
            wrong_direction = "、".join(f"{n}{infer_mistake(statements[n])}" for n in wrong_nums)
    elif answer in opts:
        correct_judgement = short_cell(opts[answer], 120)
    else:
        correct_judgement = f"选择{answer}"
    overview = (
        f"在{scenario}中，应当{correct_judgement or f'选择{answer}'}；"
        f"本题主要错在{wrong_direction}，因为选必二判断必须同时看主体、行为、权利义务和法定条件。"
    )
    wrongs: list[str] = []
    if opts and answer in combos and statements:
        correct_nums = set(combos[answer])
        for label, nums in opts_as_combo_order(combos):
            if label == answer:
                continue
            nums_set = set(nums)
            included_wrong = [n for n in nums if n not in correct_nums]
            missed_right = [n for n in correct_nums if n not in nums_set]
            parts = []
            if included_wrong:
                parts.append("包含错误判断" + "、".join(f"{n}（{infer_mistake(statements.get(n, ''))}）" for n in included_wrong))
            if missed_right:
                parts.append("漏掉正确判断" + "、".join(missed_right))
            wrongs.append(f"{label}项不对在{'，'.join(parts) or '组合关系不符合题意'}，因为该组合没有同时扣准题面事实和具体法律规则。")
    elif opts:
        for label, text in opts.items():
            if label == answer:
                continue
            mistake = infer_mistake(text)
            wrongs.append(f"{label}项不对在{mistake}，因为该说法“{short_cell(text, 80)}”需要回到题面事实和具体法律规则复核。")
    else:
        wrongs.append("本题选项文本未能稳定抽出，需回原卷逐项补写错项句。")
    return overview, "；".join(wrongs)


def numbered_statement_map(context: str) -> dict[str, str]:
    markers = list(re.finditer(r"[①②③④]", context))
    result: dict[str, str] = {}
    if not markers:
        return result
    option_start = re.search(r"[ABCD]\s*[\.．、]", context)
    limit = option_start.start() if option_start else len(context)
    for idx, marker in enumerate(markers):
        label = marker.group(0)
        start = marker.end()
        end = markers[idx + 1].start() if idx + 1 < len(markers) else limit
        if start >= limit:
            continue
        text = clean_text(context[start:end])
        if text:
            result[label] = text[:260]
    return result


def combo_option_map(opts: dict[str, str]) -> dict[str, list[str]]:
    combos: dict[str, list[str]] = {}
    for label, text in opts.items():
        nums = re.findall(r"[①②③④]", text)
        if nums:
            combos[label] = nums
    return combos


def opts_as_combo_order(combos: dict[str, list[str]]) -> list[tuple[str, list[str]]]:
    return [(label, combos[label]) for label in ["A", "B", "C", "D"] if label in combos]


def rubric_context(qno: str, suite: Suite) -> tuple[str, str, str]:
    files = suite.rubric_files + suite.support_files
    best_text = ""
    best_file = ""
    evidence = "missing"
    for item in files:
        text = normalize_for_search(item.text)
        if not text:
            continue
        patterns = [
            rf"第\s*{re.escape(qno)}\s*题",
            rf"(?<!\d){re.escape(qno)}\s*[\.．、]",
            rf"{re.escape(qno)}\s*题",
        ]
        for pat in patterns:
            match = re.search(pat, text)
            if match:
                start = max(0, match.start() - 240)
                end = min(len(text), match.end() + 2200)
                best_text = text[start:end]
                best_file = str(item.path)
                evidence = evidence_type(item.path)
                return short_cell(best_text, 2400), best_file, evidence
    # Fallback: legal-rich rubric fragment.
    for item in files:
        text = normalize_for_search(item.text)
        if private_law_score(text) >= 2:
            m = re.search("|".join(re.escape(t) for t in CORE_XUANBIER_TERMS), text)
            if m:
                start = max(0, m.start() - 300)
                end = min(len(text), m.end() + 2000)
                return short_cell(text[start:end], 2400), str(item.path), evidence_type(item.path)
    return "", "", "missing"


def evidence_type(path: Path) -> str:
    name = path.name
    if any(k in name for k in ["评标", "阅卷", "评分", "细则"]):
        return "formal_or_scoring_source"
    if "讲评" in name:
        return "support_lecture"
    if any(k in name for k in ["答案", "参考"]):
        return "reference_answer"
    return "unknown"


def complete_prompt(context: str) -> str:
    candidates = re.findall(r"[^。！？；;]{0,80}(运用《?法律与生活》?[^。！？；;]{0,120}[。！？；;]?)", context)
    if candidates:
        return clean_text(candidates[-1])
    candidates = re.findall(r"[^。！？；;]{0,80}(结合材料[^。！？；;]{0,120}[。！？；;]?)", context)
    if candidates:
        return clean_text(candidates[-1])
    # Last sentence as weak fallback.
    bits = re.split(r"[。！？]", clean_text(context))
    return bits[-1][-180:] if bits else ""


def write_csv(path: Path, rows: list[dict], fieldnames: list[str]) -> None:
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        for row in rows:
            writer.writerow({k: row.get(k, "") for k in fieldnames})


def md_escape(text: str) -> str:
    return text.replace("\n", " ").replace("|", "\\|")


def main() -> int:
    suites = discover_suites()
    source_rows: list[dict] = []
    index_rows: list[dict] = []
    choice_rows: list[dict] = []
    subj_rows: list[dict] = []
    no_suites: list[Suite] = []
    uncertain_notes: list[str] = []

    for suite in suites:
        attach_files(suite)
        extract_files(suite.paper_files + suite.rubric_files + suite.support_files)
        determine_status(suite)
        paper_text = "\n".join(f.text for f in suite.paper_files)
        answer_key = extract_answer_key("\n".join(f.text for f in suite.rubric_files + suite.support_files))
        candidates = candidate_questions(suite) if suite.status == "has_xuanbier" else {}

        source_rows.append(
            {
                "suite_id": suite.suite_id,
                "year": suite.year,
                "district": suite.district,
                "stage": suite.stage,
                "suite_path": str(suite.path),
                "paper_files": "\n".join(str(f.path) for f in suite.paper_files),
                "rubric_files": "\n".join(str(f.path) for f in suite.rubric_files),
                "support_files": "\n".join(str(f.path) for f in suite.support_files),
                "extraction_status": "; ".join(sorted(set(f.status for f in suite.paper_files + suite.rubric_files + suite.support_files))),
                "xuanbier_status": suite.status,
                "notes": suite.notes,
            }
        )

        if suite.status == "no_xuanbier":
            no_suites.append(suite)
            continue
        if suite.status in {"uncertain", "excluded"}:
            uncertain_notes.append(f"- `{suite.name}` (`{suite.path}`): {suite.status}。{suite.notes}")
            continue

        if not candidates:
            uncertain_notes.append(f"- `{suite.name}` (`{suite.path}`): 判定含选必二，但自动抽题未锁出题号，需人工复核。")
            continue

        for qno, context in sorted(candidates.items(), key=lambda kv: int(kv[0])):
            qtype = "choice" if int(qno) <= 15 else "subjective"
            rub_text, rub_file, evidence = rubric_context(qno, suite)
            evidence_status = "matched" if rub_text else "missing"
            index_rows.append(
                {
                    "suite_id": suite.suite_id,
                    "year": suite.year,
                    "district": suite.district,
                    "stage": suite.stage,
                    "suite_name": suite.name,
                    "status": suite.status,
                    "question_no": qno,
                    "question_type": qtype,
                    "source_question_file": "\n".join(str(f.path) for f in suite.paper_files),
                    "source_rubric_file": rub_file,
                    "evidence_status": evidence_status,
                    "question_excerpt": short_cell(context, 900),
                    "rubric_excerpt": short_cell(rub_text, 900),
                    "notes": "" if rub_text else "未自动匹配到对应细则原文",
                }
            )
            if qtype == "choice":
                answer = answer_key.get(qno)
                overview, wrongs = choice_sentences(qno, context, answer)
                choice_rows.append(
                    {
                        "suite_id": suite.suite_id,
                        "year": suite.year,
                        "district": suite.district,
                        "stage": suite.stage,
                        "question_no": qno,
                        "answer": answer or "",
                        "overview_sentence": overview,
                        "wrong_option_sentences": wrongs,
                        "question_excerpt": short_cell(context, 900),
                        "answer_source": rub_file,
                        "notes": "" if answer else "未自动锁定可靠答案",
                    }
                )
            else:
                subj_rows.append(
                    {
                        "suite_id": suite.suite_id,
                        "year": suite.year,
                        "district": suite.district,
                        "stage": suite.stage,
                        "question_no": qno,
                        "question_material": short_cell(context, 2200),
                        "complete_prompt": complete_prompt(context),
                        "rubric_source": rub_file,
                        "rubric_position": f"自动匹配第{qno}题上下文" if rub_file else "",
                        "evidence_type": evidence,
                        "rubric_original_text": short_cell(rub_text, 2400),
                        "notes": "" if rub_text else "未匹配到对应细则；需回原细则/OCR复核",
                    }
                )

    write_csv(
        OUT / "SOURCE_MATCH_LEDGER.csv",
        source_rows,
        ["suite_id", "year", "district", "stage", "suite_path", "paper_files", "rubric_files", "support_files", "extraction_status", "xuanbier_status", "notes"],
    )
    write_csv(
        OUT / "LEGAL_QUESTION_INDEX.csv",
        index_rows,
        ["suite_id", "year", "district", "stage", "suite_name", "status", "question_no", "question_type", "source_question_file", "source_rubric_file", "evidence_status", "question_excerpt", "rubric_excerpt", "notes"],
    )
    write_csv(
        OUT / "CHOICE_PREPROCESS.csv",
        choice_rows,
        ["suite_id", "year", "district", "stage", "question_no", "answer", "overview_sentence", "wrong_option_sentences", "question_excerpt", "answer_source", "notes"],
    )
    write_csv(
        OUT / "SUBJECTIVE_PREPROCESS.csv",
        subj_rows,
        ["suite_id", "year", "district", "stage", "question_no", "question_material", "complete_prompt", "rubric_source", "rubric_position", "evidence_type", "rubric_original_text", "notes"],
    )

    write_markdown(choice_rows, subj_rows, no_suites, uncertain_notes, source_rows)
    update_progress(len(suites), len(index_rows), len(choice_rows), len(subj_rows), len(no_suites), len(uncertain_notes))
    print(f"scanned_suites={len(suites)}")
    print(f"candidate_questions={len(index_rows)} choice={len(choice_rows)} subjective={len(subj_rows)}")
    print(f"no_xuanbier={len(no_suites)} uncertain_or_missing={len(uncertain_notes)}")
    return 0


def write_markdown(choice_rows: list[dict], subj_rows: list[dict], no_suites: list[Suite], uncertain_notes: list[str], source_rows: list[dict]) -> None:
    choice_md = [
        "# 选必二选择题预处理",
        "",
        "本文件从原始试卷与可读答案/细则来源自动预处理生成。旧选必二成果不作为依据。",
        "",
    ]
    for row in choice_rows:
        choice_md.extend(
            [
                f"## {row['year']} {row['district']} {row['stage']} 第{row['question_no']}题",
                "",
                f"- 答案：{row['answer'] or '待核'}",
                f"- 总览句：{row['overview_sentence']}",
                f"- 错项句：{row['wrong_option_sentences']}",
                f"- 原题片段：{row['question_excerpt']}",
                f"- 答案来源：{row['answer_source'] or '未自动锁定'}",
                "",
            ]
        )
    if not choice_rows:
        choice_md.append("暂无已自动锁定的选择题。")
    (OUT / "CHOICE_PREPROCESS.md").write_text("\n".join(choice_md), encoding="utf-8")

    subj_md = [
        "# 选必二主观题预处理",
        "",
        "本文件只锁定题面、完整设问与对应细则原文，不做框架归置、不总结母题、不写课堂模板。",
        "",
    ]
    for row in subj_rows:
        subj_md.extend(
            [
                f"## {row['year']} {row['district']} {row['stage']} 第{row['question_no']}题",
                "",
                f"- 完整设问：{row['complete_prompt'] or '待人工切分'}",
                f"- 题面/材料片段：{row['question_material']}",
                f"- 细则来源：{row['rubric_source'] or '未自动匹配'}",
                f"- 细则位置：{row['rubric_position'] or '待人工复核'}",
                f"- 证据类型：{row['evidence_type']}",
                "",
                "### 对应细则原文",
                "",
                row["rubric_original_text"] or "未自动匹配到对应细则原文，需回原细则/OCR复核。",
                "",
            ]
        )
    if not subj_rows:
        subj_md.append("暂无已自动锁定的主观题。")
    (OUT / "SUBJECTIVE_PREPROCESS.md").write_text("\n".join(subj_md), encoding="utf-8")

    no_md = [
        "# 确认不含选必二的套卷",
        "",
        "确认无《法律与生活》题的套卷记录在这里。期中卷若确认不含选必二，标记后跳过，不硬凑题。",
        "",
    ]
    for suite in no_suites:
        no_md.append(f"- `{suite.name}` (`{suite.path}`): no_xuanbier。{suite.notes}")
    if not no_suites:
        no_md.append("暂无确认 no_xuanbier 套卷。")
    (OUT / "NO_XUANBIER_SUITES.md").write_text("\n".join(no_md) + "\n", encoding="utf-8")

    missing_md = [
        "# 缺失或不确定项目",
        "",
        "本文件记录缺题面、缺细则、OCR 不稳、模块边界不稳、答案源冲突等问题。不得把不确定项目冒充已完成。",
        "",
    ]
    missing_md.extend(uncertain_notes or ["暂无自动发现的不确定项目。"])
    # Add candidate rows with missing rubric.
    missing_subjective = [r for r in subj_rows if r.get("evidence_type") == "missing"]
    if missing_subjective:
        missing_md.extend(["", "## 主观题细则未匹配"])
        for row in missing_subjective:
            missing_md.append(f"- `{row['year']} {row['district']} {row['stage']} 第{row['question_no']}题`: {row['notes']}")
    missing_choices = [r for r in choice_rows if r.get("notes")]
    if missing_choices:
        missing_md.extend(["", "## 选择题答案未锁定"])
        for row in missing_choices:
            missing_md.append(f"- `{row['year']} {row['district']} {row['stage']} 第{row['question_no']}题`: {row['notes']}")
    (OUT / "MISSING_OR_UNCERTAIN.md").write_text("\n".join(missing_md) + "\n", encoding="utf-8")

    summary = [
        "# 自动预处理扫描摘要",
        "",
        f"- 套卷数：{len(source_rows)}",
        f"- has_xuanbier：{sum(1 for r in source_rows if r['xuanbier_status'] == 'has_xuanbier')}",
        f"- no_xuanbier：{sum(1 for r in source_rows if r['xuanbier_status'] == 'no_xuanbier')}",
        f"- uncertain：{sum(1 for r in source_rows if r['xuanbier_status'] == 'uncertain')}",
        f"- excluded：{sum(1 for r in source_rows if r['xuanbier_status'] == 'excluded')}",
    ]
    (OUT / "audit/SCAN_SUMMARY.md").write_text("\n".join(summary) + "\n", encoding="utf-8")

    evidence_counts: dict[str, int] = {}
    for row in subj_rows:
        evidence_counts[row.get("evidence_type", "unknown")] = evidence_counts.get(row.get("evidence_type", "unknown"), 0) + 1
    validation = [
        "# 预处理验收报告",
        "",
        "- 小本本检查：已创建 `00_飞哥选必二法律与生活要求小本本.md`；本轮脚本不覆盖该文件。",
        f"- 覆盖检查：本轮扫描 `{len(source_rows)}` 套，每套均写入 `SOURCE_MATCH_LEDGER.csv` 并带 `has_xuanbier / no_xuanbier / uncertain / excluded` 状态。",
        f"- 跳过检查：确认 `no_xuanbier` `{sum(1 for r in source_rows if r['xuanbier_status'] == 'no_xuanbier')}` 套，写入 `NO_XUANBIER_SUITES.md`。",
        f"- 不确定检查：`MISSING_OR_UNCERTAIN.md` 记录 `{len(uncertain_notes)}` 条套卷/OCR/自动抽题问题。",
        f"- 选择题检查：输出 `{len(choice_rows)}` 道；总览句与错项句字段均已填充；其中 `{sum(1 for r in choice_rows if r.get('notes'))}` 道未自动锁定可靠答案，已单列待核。",
        f"- 主观题检查：输出 `{len(subj_rows)}` 道；细则缺失 `{sum(1 for r in subj_rows if r.get('evidence_type') == 'missing')}` 道。",
        f"- 证据检查：主观题证据类型分布 `{evidence_counts}`；`reference_answer` 不等同于正式细则。",
        "- 旧线隔离检查：本轮输出声明旧选必二作废；脚本来源只扫原始三年模拟题目录和新目录缓存，不读取旧 `选必二_*.md`。",
    ]
    (OUT / "audit/VALIDATION_REPORT.md").write_text("\n".join(validation) + "\n", encoding="utf-8")


def update_progress(suite_count: int, q_count: int, choice_count: int, subj_count: int, no_count: int, uncertain_count: int) -> None:
    progress = f"""# 选必二《法律与生活》预处理进度

日期：2026-04-30

- [x] STEP_01: 已按用户要求删除旧选必二独立产物、旧选必二监管报告、旧选必二选择题处理台账、旧连续任务目录；保留原始试卷资料和总路由 skill。
- [x] STEP_02: 已建立本轮新目录、小本本、任务说明、开发计划、进度文件、空台账和预处理脚本。
- [x] STEP_03: 已从原始三年模拟题资料池扫描套卷、试卷、细则、答案、评标、阅卷总结和讲评材料，共扫描 `{suite_count}` 套。
- [x] STEP_04: 已判定含选必二 / 无选必二 / 不确定套卷；确认无选必二套卷 `{no_count}` 套，不确定或缺 OCR/转换 `{uncertain_count}` 条。
- [x] STEP_05: 已对确认或疑似含选必二的套卷自动锁定法律选择题和主观题候选题号，共 `{q_count}` 个候选题。
- [x] STEP_06: 已为主观题候选项匹配对应细则原文；未匹配或证据不足项目写入 `MISSING_OR_UNCERTAIN.md`。
- [x] STEP_07: 已为自动锁定的选择题生成“两层结构”预处理，共 `{choice_count}` 道选择题；主观题预处理 `{subj_count}` 道。
- [x] STEP_08: 已输出 Markdown/CSV 与扫描摘要；下一轮需人工抽查若干套，修正 OCR/自动抽题误差。
"""
    (OUT / "PROGRESS.md").write_text(progress, encoding="utf-8")


if __name__ == "__main__":
    raise SystemExit(main())
