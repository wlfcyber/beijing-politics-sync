#!/usr/bin/env python3
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import os
import re
import subprocess
import sys
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Iterable


RUN_DIR = Path(__file__).resolve().parents[1]
DESKTOP = Path("/Users/wanglifei/Desktop")
CONTROL = RUN_DIR / "00_control"
INVENTORY = RUN_DIR / "01_inventory"
TEXT_DIR = RUN_DIR / "02_extracted_text"
DRAFT_DIR = RUN_DIR / "04_working_draft"
OUTPUT_DIR = RUN_DIR / "05_output"

SUPPORTED_EXTS = {
    ".docx",
    ".doc",
    ".pdf",
    ".txt",
    ".md",
    ".pptx",
    ".ppt",
    ".jpg",
    ".jpeg",
    ".png",
    ".heic",
}

RAW_ROOTS = [
    DESKTOP / "2024模拟题",
    DESKTOP / "2025模拟题",
    DESKTOP / "2026模拟题",
]

SUPPLEMENTAL_SOURCE_HINTS = [
    "模拟题资料",
    "期中讲评",
    "期末讲评",
    "一模讲评",
    "二模讲评",
]

EXCLUDE_CONTENT_KEYWORDS = [
    "汇编",
    "题包",
    "整理稿",
    "宝典",
    "最终",
    "终极",
    "定稿",
    "修订",
    "校验",
    "报告",
    "框架",
    "术语",
    "频次",
    "候选",
    "学生版",
    "审核稿",
    "Claude",
    "Codex",
    "飞哥政治庄园",
    "北京高考政治/选必",
    "北京高考政治/必修",
    "raw_exam_subjective_compilation_2026-06-02",
    "经济与社会",
    "考题-wlf",
    "extracted_text",
    "rendered_pages",
    "page_",
    "_strip",
]

CLASSIFICATION_KEYWORDS = ["试题分类", "分类汇编", "按模块"]
QUESTION_KEYWORDS = ["试题", "试卷", "政治", "高三", "一模", "二模", "期中", "期末", "统练", "练习", "检测"]
RUBRIC_KEYWORDS = ["细则", "评分", "答案", "参考", "评标", "阅卷", "讲评", "解析"]
LEGAL_KEYWORDS = [
    "法律与生活",
    "民法典",
    "合同",
    "侵权",
    "人格权",
    "知识产权",
    "劳动",
    "劳动合同",
    "消费者",
    "继承",
    "婚姻",
    "仲裁",
    "诉讼",
    "调解",
    "司法确认",
    "相邻关系",
    "不正当竞争",
    "格式条款",
    "著作权",
    "商标权",
    "专利权",
    "所有权",
    "用益物权",
    "担保物权",
]


@dataclass
class SourceRow:
    source_id: str
    path: str
    year: str
    district_or_exam: str
    paper_type: str
    file_role: str
    file_ext: str
    include_status: str
    reason: str
    conversion_status: str
    text_path: str
    sha256: str
    size_bytes: int
    mtime: str


def is_under(path: Path, root: Path) -> bool:
    try:
        path.resolve().relative_to(root.resolve())
        return True
    except ValueError:
        return False


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def short_id(path: Path) -> str:
    digest = hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:10]
    name = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff]+", "_", path.stem)[:42].strip("_")
    return f"SRC_{digest}_{name}"


def parse_year(text: str) -> str:
    m = re.search(r"20(24|25|26)", text)
    return m.group(0) if m else ""


def parse_paper_type(text: str) -> str:
    matches: list[tuple[int, str]] = []
    for token in ["一模", "二模", "期中", "期末", "适应性", "高考", "联考", "统练", "练习", "检测"]:
        pos = text.rfind(token)
        if pos >= 0:
            matches.append((pos, token))
    if matches:
        return sorted(matches)[-1][1]
    return ""


def parse_district(text: str) -> str:
    districts = [
        "海淀",
        "西城",
        "东城",
        "朝阳",
        "丰台",
        "石景山",
        "门头沟",
        "房山",
        "通州",
        "顺义",
        "昌平",
        "大兴",
        "怀柔",
        "平谷",
        "密云",
        "延庆",
        "北京",
    ]
    hits = [d for d in districts if d in text]
    if hits:
        return "、".join(dict.fromkeys(hits))
    return ""


def classify_role(path: Path) -> str:
    s = str(path)
    if any(k in s for k in CLASSIFICATION_KEYWORDS):
        return "classification_compilation"
    if any(k in s for k in ["细则", "评分", "评标", "阅卷"]):
        return "rubric_or_marking"
    if any(k in s for k in ["答案", "参考"]):
        return "answer_reference"
    if "讲评" in s or "解析" in s:
        return "lecture_or_explanation"
    if any(k in s for k in QUESTION_KEYWORDS):
        return "question_paper"
    return "unknown"


def include_decision(path: Path) -> tuple[str, str]:
    s = str(path)
    if path.name.startswith("~$") or "/." in s:
        return "exclude", "temporary_or_hidden"
    if "2026石景山期末" in s:
        return "exclude", "excluded_by_xuanbier_notebook_no_usable_rubric"
    if any(k in s for k in CLASSIFICATION_KEYWORDS):
        return "exclude", "classification_compilation_not_independent_suite"
    if any(k in s for k in EXCLUDE_CONTENT_KEYWORDS):
        return "exclude", "old_or_generated_artifact"
    if any(is_under(path, root) for root in RAW_ROOTS):
        return "include", "under_raw_year_source_root"
    if parse_year(s) and any(h in s for h in SUPPLEMENTAL_SOURCE_HINTS):
        if "北京高考政治/" not in s and "压缩包与交接包" not in s and "Codex" not in s:
            return "include", "desktop_supplemental_source_hint"
    return "exclude", "not_exam_source_scope"


def discover_files() -> list[SourceRow]:
    rows: list[SourceRow] = []
    for path in DESKTOP.rglob("*"):
        if not path.is_file():
            continue
        ext = path.suffix.lower()
        if ext not in SUPPORTED_EXTS:
            continue
        s = str(path)
        if not (parse_year(s) or any(is_under(path, root) for root in RAW_ROOTS)):
            continue
        status, reason = include_decision(path)
        stat = path.stat()
        rows.append(
            SourceRow(
                source_id=short_id(path),
                path=s,
                year=parse_year(s),
                district_or_exam=parse_district(s),
                paper_type=parse_paper_type(s),
                file_role=classify_role(path),
                file_ext=ext,
                include_status=status,
                reason=reason,
                conversion_status="pending" if status == "include" else "not_converted",
                text_path="",
                sha256=sha256_file(path) if status == "include" else "",
                size_bytes=stat.st_size,
                mtime=str(int(stat.st_mtime)),
            )
        )
    rows.sort(key=lambda r: (r.include_status != "include", r.year, r.path))
    return rows


def write_csv(path: Path, rows: Iterable[dict], fieldnames: list[str]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        for row in rows:
            writer.writerow(row)


def load_source_rows() -> list[SourceRow]:
    path = CONTROL / "SOURCE_LEDGER.csv"
    with path.open("r", encoding="utf-8-sig", newline="") as f:
        reader = csv.DictReader(f)
        return [SourceRow(**{**row, "size_bytes": int(row["size_bytes"]) if row.get("size_bytes") else 0}) for row in reader]


def text_from_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.rstrip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def text_from_pdf(path: Path) -> str:
    chunks: list[str] = []
    try:
        import pdfplumber

        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    chunks.append(f"\n[PDF_PAGE {i}]\n{text}")
    except Exception:
        chunks = []
    if chunks:
        return "\n".join(chunks)
    try:
        import fitz

        doc = fitz.open(str(path))
        for i, page in enumerate(doc, 1):
            text = page.get_text("text") or ""
            if text.strip():
                chunks.append(f"\n[PDF_PAGE {i}]\n{text}")
    except Exception:
        pass
    return "\n".join(chunks)


def text_from_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        slide_text.append(" | ".join(cells))
        if slide_text:
            chunks.append(f"\n[PPT_SLIDE {i}]\n" + "\n".join(slide_text))
    return "\n".join(chunks)


def text_from_text_file(path: Path) -> str:
    for enc in ["utf-8-sig", "utf-8", "gb18030", "utf-16"]:
        try:
            return path.read_text(encoding=enc)
        except Exception:
            continue
    return path.read_bytes().decode("utf-8", errors="replace")


def text_from_legacy_office(path: Path) -> str:
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            check=False,
            capture_output=True,
            text=True,
            timeout=120,
        )
    except Exception as exc:
        return f"[CONVERSION_ERROR] {type(exc).__name__}: {exc}"
    return result.stdout or result.stderr or ""


def text_from_ocr(path: Path) -> str:
    ocr = Path("/Users/wanglifei/.local/bin/ocr-vision")
    if not ocr.exists():
        return ""
    try:
        result = subprocess.run(
            [str(ocr), str(path)],
            check=False,
            capture_output=True,
            text=True,
            timeout=180,
        )
    except Exception:
        return ""
    return (result.stdout or "").strip()


def extract_one(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    try:
        if ext == ".docx":
            text = text_from_docx(path)
        elif ext in {".doc", ".ppt"}:
            text = text_from_legacy_office(path)
        elif ext == ".pdf":
            text = text_from_pdf(path)
            if len(text.strip()) < 50:
                ocr_text = text_from_ocr(path)
                if ocr_text:
                    return ocr_text, "ocr"
        elif ext == ".pptx":
            text = text_from_pptx(path)
        elif ext in {".txt", ".md"}:
            text = text_from_text_file(path)
        elif ext in {".jpg", ".jpeg", ".png", ".heic"}:
            text = text_from_ocr(path)
            return text, "ocr" if text else "blocked_conversion"
        else:
            return "", "unsupported_ext"
    except Exception as exc:
        return f"[CONVERSION_ERROR] {type(exc).__name__}: {exc}", "conversion_error"
    return text, "converted" if text.strip() else "empty_text"


def inventory_cmd() -> None:
    rows = discover_files()
    fieldnames = list(asdict(rows[0]).keys()) if rows else [f.name for f in SourceRow.__dataclass_fields__.values()]
    write_csv(CONTROL / "SOURCE_LEDGER.csv", (asdict(r) for r in rows), fieldnames)
    write_csv(INVENTORY / "desktop_candidate_sources.csv", (asdict(r) for r in rows), fieldnames)

    included = [r for r in rows if r.include_status == "include"]
    excluded = [r for r in rows if r.include_status != "include"]
    md = ["# Desktop Candidate Source Inventory", "", "## Included Candidate Sources", ""]
    md.append("| source_id | year | district/exam | paper_type | role | file | reason |")
    md.append("| --- | --- | --- | --- | --- | --- | --- |")
    for r in included:
        md.append(
            f"| {r.source_id} | {r.year} | {r.district_or_exam} | {r.paper_type} | {r.file_role} | {Path(r.path).name} | {r.reason} |"
        )
    md.extend(["", "## Excluded/Not Evidence", ""])
    md.append("| year | role | file | reason |")
    md.append("| --- | --- | --- | --- |")
    for r in excluded:
        md.append(f"| {r.year} | {r.file_role} | {Path(r.path).name} | {r.reason} |")
    (INVENTORY / "desktop_candidate_sources.md").write_text("\n".join(md) + "\n", encoding="utf-8")

    print("候选原始来源清单（include）")
    print("source_id\tyear\tdistrict/exam\tpaper_type\trole\tfile")
    for r in included:
        print(f"{r.source_id}\t{r.year}\t{r.district_or_exam}\t{r.paper_type}\t{r.file_role}\t{Path(r.path).name}")
    print(f"\ninclude={len(included)} exclude_or_not_evidence={len(excluded)}")
    print(f"wrote {INVENTORY / 'desktop_candidate_sources.md'}")


def extract_text_cmd() -> None:
    rows = load_source_rows()
    updated: list[SourceRow] = []
    for row in rows:
        if row.include_status != "include":
            updated.append(row)
            continue
        path = Path(row.path)
        out = TEXT_DIR / f"{row.source_id}.txt"
        if out.exists() and out.stat().st_size > 0:
            row.conversion_status = row.conversion_status if row.conversion_status != "pending" else "converted_cached"
            row.text_path = str(out)
            updated.append(row)
            continue
        text, status = extract_one(path)
        out.write_text(text, encoding="utf-8")
        row.conversion_status = status
        row.text_path = str(out)
        updated.append(row)
        print(f"{row.source_id}\t{status}\t{path.name}")
    write_csv(CONTROL / "SOURCE_LEDGER.csv", (asdict(r) for r in updated), list(asdict(updated[0]).keys()))
    write_csv(INVENTORY / "desktop_candidate_sources.csv", (asdict(r) for r in updated), list(asdict(updated[0]).keys()))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("command", choices=["inventory", "extract-text"])
    args = parser.parse_args()
    if args.command == "inventory":
        inventory_cmd()
    elif args.command == "extract-text":
        extract_text_cmd()


if __name__ == "__main__":
    main()
